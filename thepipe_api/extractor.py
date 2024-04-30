import base64
from concurrent.futures import ThreadPoolExecutor
from io import BytesIO
import math
import re
from typing import Dict, List, Optional
import glob
import os
import tempfile
from urllib.parse import urlparse
import zipfile
from PIL import Image
import requests
import json
from .core import Chunk, print_status, SourceTypes, create_chunks_from_messages, API_URL
import tempfile
import mimetypes
import dotenv
from datetime import datetime
from json import JSONEncoder
dotenv.load_dotenv()

class JSONDateEncoder(JSONEncoder):
    def default(self, obj):
        if isinstance(obj, datetime):
            return obj.isoformat()
        return JSONEncoder.default(self, obj)

FILES_TO_IGNORE = {'LICENSE', 'package-lock.json', '.gitignore', '.bin', '.pyc', '.pyo', '.exe', '.bat', '.dll', '.obj', '.o', '.a', '.lib', '.so', '.dylib', '.ncb', '.sdf', '.suo', '.pdb', '.idb', '.pyd', '.ipynb_checkpoints', '.npy', '.pth'} # Files to ignore, please feel free to customize!
CODE_EXTENSIONS = {'.h', '.json', '.js', '.jsx', '.ts', '.tsx',  '.cs', '.java', '.html', '.css', '.xml', '.yaml', '.xaml', '.sh'} # Plaintext files that should not be compressed with LLMLingua
CTAGS_CODE_EXTENSIONS = {'.c', '.cpp', '.py'} # code files that work with ctags
PLAINTEXT_EXTENSIONS = {'.txt', '.md', '.rtf', '.ino', '.ini', '.cfg', '.log'}
IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png'}
SPREADSHEET_EXTENSIONS = {'.csv', '.xls', '.xlsx'}
DOCUMENT_EXTENSIONS = {'.pdf', '.docx', '.pptx'}
VIDEO_EXTENSIONS = {'.mp4', '.avi', '.mov', '.mkv'}
AUDIO_EXTENSIONS = {'.mp3', '.wav'}
OTHER_EXTENSIONS = {'.zip', '.ipynb'}
KNOWN_EXTENSIONS = IMAGE_EXTENSIONS.union(CODE_EXTENSIONS).union(CTAGS_CODE_EXTENSIONS).union(PLAINTEXT_EXTENSIONS).union(IMAGE_EXTENSIONS).union(SPREADSHEET_EXTENSIONS).union(DOCUMENT_EXTENSIONS).union(OTHER_EXTENSIONS)
GITHUB_TOKEN: str = os.getenv("GITHUB_TOKEN")
THEPIPE_API_KEY: str = os.getenv("THEPIPE_API_KEY")

def extract_from_source(source: str, match: Optional[str] = None, ignore: Optional[str] = None, limit: int = None, verbose: bool = False, ai_extraction: bool = False, text_only: bool = False, local: bool = True) -> List[Chunk]:
    source_type = detect_type(source)
    if source_type is None:
        raise ValueError(f"Could not detect source type for {source}.")
    if verbose: print_status(f"Extracting from {source_type.value}", status='info')
    if source_type == SourceTypes.DIR or source == '.' or source == './':
        if source == '.' or source == './':
            source = os.getcwd()
        return extract_from_directory(dir_path=source, match=match, ignore=ignore, verbose=verbose, ai_extraction=ai_extraction, text_only=text_only, limit=limit, local=local)
    elif source_type == SourceTypes.GITHUB:
        return extract_github(github_url=source, file_path='', match=match, ignore=ignore, text_only=text_only, verbose=verbose, ai_extraction=ai_extraction, branch='master')
    elif source_type == SourceTypes.YOUTUBE_VIDEO:
        return extract_youtube(youtube_url=source, text_only=text_only, verbose=verbose)
    elif source_type == SourceTypes.URL:
        return extract_url(url=source, text_only=text_only, local=local)
    return extract_from_file(file_path=source, source_type=source_type, verbose=verbose, ai_extraction=ai_extraction, text_only=text_only, local=local)

def extract_from_file(file_path: str, source_type: str, verbose: bool = False, ai_extraction: bool = False, text_only: bool = False, local: bool = True, limit: int = None) -> List[Chunk]:
    if not local:
        try:
            with open(file_path, 'rb') as f:
                response = requests.post(
                    url=API_URL,
                    files={'file': (file_path, f)},
                    data={'api_key': THEPIPE_API_KEY, 'ai_extraction': ai_extraction, 'text_only': text_only, 'limit': limit}
                )
        except Exception as e:
            raise ValueError(f"Failed to extract from {file_path}. This may mean our backend couldn't handle this request. Exception: {e}.")
        try:
            response = response.json()
        except json.JSONDecodeError:
            raise ValueError(f"Failed to extract from {file_path}. This may mean our backend couldn't handle this request. Response: {response}.")
        if 'error' in response:
            raise ValueError(f"{response['error']}")
        chunks = create_chunks_from_messages(response['messages'])
        for c in chunks:
            c.path = file_path
            c.source_type = source_type
        return chunks
    try:    
        if source_type == SourceTypes.PDF:
            extraction = extract_pdf(file_path=file_path, ai_extraction=ai_extraction, text_only=text_only, verbose=verbose)
        elif source_type == SourceTypes.DOCX:
            extraction = extract_docx(file_path=file_path, verbose=verbose, text_only=text_only)
        elif source_type == SourceTypes.PPTX:
            extraction = extract_pptx(file_path=file_path, verbose=verbose, text_only=text_only)
        elif source_type == SourceTypes.IMAGE:
            extraction = [extract_image(file_path=file_path, text_only=text_only)]
        elif source_type == SourceTypes.SPREADSHEET:
            extraction = extract_spreadsheet(file_path=file_path)
        elif source_type == SourceTypes.PLAINTEXT:
            extraction = [extract_plaintext(file_path=file_path)]
        elif source_type == SourceTypes.UNCOMPRESSIBLE_CODE:
            extraction = [extract_plaintext(file_path=file_path)]
            extraction = [Chunk(path=e.path, text=e.text, image=None, source_type=SourceTypes.UNCOMPRESSIBLE_CODE) for e in extraction] # change types to code
        elif source_type == SourceTypes.COMPRESSIBLE_CODE:
            extraction = [extract_plaintext(file_path=file_path)]
            extraction = [Chunk(path=e.path, text=e.text, image=None, source_type=SourceTypes.COMPRESSIBLE_CODE) for e in extraction]
        elif source_type == SourceTypes.IPYNB:
            extraction = extract_from_ipynb(file_path=file_path, verbose=verbose, ai_extraction=ai_extraction, text_only=text_only)
        elif source_type == SourceTypes.ZIP:
            extraction = extract_zip(file_path=file_path, verbose=verbose, ai_extraction=ai_extraction, text_only=text_only)
        elif source_type == SourceTypes.VIDEO:
            extraction = extract_video(file_path=file_path, verbose=verbose, text_only=text_only)
        elif source_type == SourceTypes.AUDIO:
            extraction = extract_audio(file_path=file_path, verbose=verbose)
        else:
            extraction = [extract_plaintext(file_path=file_path)]
        if verbose: print_status(f"Extracted from {file_path}", status='success')
        return extraction
    except Exception as e:
        if verbose: print_status(f"Failed to extract from {file_path}: {e}", status='error')
        return [Chunk(path=file_path)]

def detect_type(source: str) -> Optional[SourceTypes]:
    if source.startswith("https://www.youtube.com") or source.startswith("https://youtube.com"):
        return SourceTypes.YOUTUBE_VIDEO
    elif source.startswith("https://github.com") or source.startswith("https://www.github.com"):
        return SourceTypes.GITHUB
    elif source.startswith("http") or source.startswith("ftp."):
        return SourceTypes.URL
    elif source.endswith(".zip"):
        return SourceTypes.ZIP
    elif os.path.isdir(source) or source == '.' or source == './':
        return SourceTypes.DIR
    # try splitting the source into a filename and extension
    _, extension = os.path.splitext(source)
    # if that fails, try to detect the file type using Magika
    if (not extension) or (extension not in KNOWN_EXTENSIONS):
        from magika import Magika # import only if needed
        magika = Magika()
        try:
            with open(source, 'rb') as file:
                result = magika.identify_bytes(file.read())
        except Exception as e:
            return None
        mimetype = result.output.mime_type
        extension = mimetypes.guess_extension(mimetype, strict=False)
    if not extension:
        return None
    # Map the detected extension to the corresponding SourceType
    if extension in CODE_EXTENSIONS:
        return SourceTypes.UNCOMPRESSIBLE_CODE
    elif extension in CTAGS_CODE_EXTENSIONS:
        return SourceTypes.COMPRESSIBLE_CODE
    elif extension in IMAGE_EXTENSIONS:
        return SourceTypes.IMAGE
    elif extension in SPREADSHEET_EXTENSIONS:
        return SourceTypes.SPREADSHEET
    elif extension == '.pdf':
        return SourceTypes.PDF
    elif extension == '.ipynb':
        return SourceTypes.IPYNB
    elif extension == '.docx':
        return SourceTypes.DOCX
    elif extension == '.pptx':
        return SourceTypes.PPTX
    elif extension == '.zip':
        return SourceTypes.ZIP
    elif extension in VIDEO_EXTENSIONS:
        return SourceTypes.VIDEO
    elif extension in AUDIO_EXTENSIONS:
        return SourceTypes.AUDIO
    elif extension in PLAINTEXT_EXTENSIONS:
        return SourceTypes.PLAINTEXT
    return None

def extract_plaintext(file_path: str) -> List[Chunk]:
    with open(file_path, 'r', encoding='utf-8') as file:
        text = file.read()
    return Chunk(path=file_path, text=text, image=None, source_type=SourceTypes.PLAINTEXT)

def should_ignore(file_path: str, ignore: Optional[str] = None) -> bool:
    if ignore is not None and re.search(ignore, file_path, re.IGNORECASE):
        return True
    if any(file_path.endswith(extension) for extension in FILES_TO_IGNORE):
        return True
    if file_path.startswith('.'):
        return True
    if any(x in file_path for x in ['node_modules', 'venv', '.git', '.vscode', '__pycache__']):
        return True
    if not os.path.isfile(file_path):
        return True
    return False

def extract_from_directory(dir_path: str, match: Optional[str] = None, ignore: Optional[str] = None, verbose: bool = False, ai_extraction: bool = False, text_only: bool = False, limit: int = None, local: bool = True) -> List[Chunk]:
    all_files = glob.glob(dir_path + "/**/*", recursive=True)
    matched_files = [file for file in all_files if re.search(match, file, re.IGNORECASE)] if match else all_files
    file_paths = [file for file in matched_files if not should_ignore(file, ignore)]
    contents = []
    with ThreadPoolExecutor() as executor:
        results = executor.map(lambda file_path: extract_from_source(source=file_path, match=match, ignore=ignore, verbose=verbose, ai_extraction=ai_extraction, text_only=text_only, limit=limit, local=local), file_paths)
        for result in results:
            contents += result
    return contents

def extract_zip(file_path: str, match: Optional[str] = None, ignore: Optional[str] = None, verbose: bool = False, ai_extraction: bool = False, text_only: bool = False) -> List[Chunk]:
    extracted_files = []
    with tempfile.TemporaryDirectory() as temp_dir:
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        extracted_files = extract_from_directory(dir_path=temp_dir, match=match, ignore=ignore, verbose=verbose, ai_extraction=ai_extraction, text_only=text_only)
    return extracted_files

def extract_pdf(file_path: str, ai_extraction: bool = False, text_only: bool = False, verbose: bool = False, limit: int = None) -> List[Chunk]:
    chunks = []
    if ai_extraction:
        with open(file_path, "rb") as f:
            response = requests.post(
                url=API_URL,
                files={'file': (file_path, f)},
                data={'api_key': THEPIPE_API_KEY, 'ai_extraction': ai_extraction, 'text_only': text_only, 'limit': limit}
            )
        try:
            response_json = response.json()
        except json.JSONDecodeError:
            raise ValueError(f"Our backend likely couldn't handle this request. This can happen with large content such as videos, streams, or very large files/websites. Re")
        if 'error' in response_json:
            raise ValueError(f"{response_json['error']}")
        messages = response_json['messages']
        chunks = create_chunks_from_messages(messages)
    else:
        import fitz
        # extract text and images of each page from the PDF
        with open(file_path, 'rb') as file:
            doc = fitz.open(file_path)
            for page in doc:
                text = page.get_text()
                if text_only:
                    chunks.append(Chunk(path=file_path, text=text, image=None, source_type=SourceTypes.PDF))
                else:
                    pix = page.get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    chunks.append(Chunk(path=file_path, text=text, image=img,source_type=SourceTypes.PDF))
            doc.close()
    return chunks

def extract_images_from_markdown(text: str) -> List[Chunk]:
    image_urls = re.findall(r"!\[.*?\]\((.*?)\)", text)
    images = []
    for url in image_urls:
        extension = os.path.splitext(urlparse(url).path)[1]
        if extension in ['.jpg', '.jpeg', '.png']:
            img = Image.open(requests.get(url, stream=True).raw)
        else:
            # ignore incompatible image extractions
            continue
        # strip url from text
        text = text.replace(f"![]({url})", "")
        images.append(Chunk(path=url, text=None, image=img, source_type=SourceTypes.IMAGE))
    return text, images

def extract_image(file_path: str, text_only: bool = False) -> Chunk:
    img = Image.open(file_path)
    img.load() # needed to close the file
    if text_only:
        import pytesseract # import only if needed
        text = pytesseract.image_to_string(img)
        return Chunk(path=file_path, text=text, image=None, source_type=SourceTypes.IMAGE)
    else:
        return Chunk(path=file_path, text=None, image=img, source_type=SourceTypes.IMAGE)
    
def extract_spreadsheet(file_path: str) -> List[Chunk]:
    import pandas as pd # import only if needed
    if file_path.endswith(".csv"):
        df = pd.read_csv(file_path)
    elif file_path.endswith(".xls") or file_path.endswith(".xlsx"):
        df = pd.read_excel(file_path)
    dicts = df.to_dict(orient='records')
    chunks = []
    for item in dicts:
        item_json = json.dumps(item, indent=4, cls=JSONDateEncoder)
        chunks.append(Chunk(path=file_path, text=item_json, image=None, source_type=SourceTypes.SPREADSHEET))
    return chunks
    
def extract_url(url: str, text_only: bool = False, local: bool = True, limit: int = None) -> List[Chunk]:
    if not local:
        try:
            response = requests.post(
                url=API_URL,
                data={'url': url, 'api_key': THEPIPE_API_KEY, 'text_only': text_only, 'limit': limit}
            )
        except Exception as e:
            raise ValueError(f"Failed to extract from URL. This may mean our backend couldn't handle this request. Exception: {e}.")
        try:
            response = response.json()
        except json.JSONDecodeError:
            raise ValueError(f"Failed to extract from URL. This may mean our backend couldn't handle this request. Response: {response}.")
        if 'error' in response:
            raise ValueError(f"{response['error']}")
        chunks = create_chunks_from_messages(response['messages'])
        return chunks
    chunks = []
    _, extension = os.path.splitext(urlparse(url).path)
    if extension is not None and extension in KNOWN_EXTENSIONS:
        # if url has a file extension, download and extract into tempfile
        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = os.path.join(temp_dir, os.path.basename(url))
            response = requests.get(url)
            with open(file_path, 'wb') as file:
                file.write(response.content)
            chunks = extract_from_source(source=file_path, text_only=text_only, local=local)
    else:
        # use playwright to extract text and images from the URL
        from playwright.sync_api import sync_playwright # import only if needed
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page()
            page.goto(url)
            if not text_only:
                # Get the viewport size and document size to scroll
                viewport_height = page.viewport_size['height']
                total_height = page.evaluate("document.body.scrollHeight")
                # Scroll to the bottom of the page and take screenshots
                current_scroll_position = 0
                scrolldowns, max_scrolldowns = 0, 10 # in case of infinite scroll
                while current_scroll_position < total_height and scrolldowns < max_scrolldowns:
                    page.wait_for_timeout(100) # wait for dynamic content to load
                    screenshot = page.screenshot()
                    img = Image.open(BytesIO(screenshot))
                    chunks.append(Chunk(path=url, text=None, image=img, source_type=SourceTypes.URL))
                    current_scroll_position += viewport_height
                    page.evaluate(f"window.scrollTo(0, {current_scroll_position})")
                    scrolldowns += 1
            text = page.inner_text('body')
            if text:
                chunks.append(Chunk(path=url, text=text, image=None, source_type=SourceTypes.URL))
            browser.close()
    if not chunks:
        raise ValueError("No content extracted from URL.")
    return chunks

def extract_video(file_path: str, verbose: bool = False, text_only: bool = False) -> List[Chunk]:
    from moviepy.editor import VideoFileClip # import only if needed
    import whisper # import only if needed
    model = whisper.load_model("base")
    video = VideoFileClip(file_path)
    chunk_duration = 60
    num_chunks = math.ceil(video.duration / chunk_duration)
    chunks = []
    for i in range(num_chunks):
        # calculate start and end time for the current chunk
        start_time = i * chunk_duration
        end_time = start_time + chunk_duration
        if end_time > video.duration:
            end_time = video.duration
        # extract frame at the middle of the chunk
        frame_time = (start_time + end_time) / 2
        frame = video.get_frame(frame_time)
        image = Image.fromarray(frame)
        # extract and transcribe audio for the current chunk
        audio_path = os.path.join(tempfile.gettempdir(), f"temp_audio_{i}.wav")
        audio = video.subclip(start_time, end_time).audio
        if audio is None:
            transcription = None
        else:
            audio.write_audiofile(audio_path, codec='pcm_s16le')
            result = model.transcribe(audio_path, verbose=verbose)
            transcription = result['text']
            os.remove(audio_path)
        # add chunk
        if not text_only:
            chunks.append(Chunk(path=file_path, text=transcription, image=image, source_type=SourceTypes.VIDEO))
        else:
            chunks.append(Chunk(path=file_path, text=transcription, image=None, source_type=SourceTypes.VIDEO))
    return chunks

def extract_youtube(youtube_url: str, text_only: bool = False, verbose: bool = False) -> List[Chunk]:
    from pytube import YouTube # import only if needed
    temp_dir = "youtube_temp"
    filename = "temp_video.mp4"
    yt = YouTube(youtube_url)
    stream = yt.streams.filter(progressive=True, file_extension='mp4').first()
    stream.download(temp_dir, filename=filename)
    video_path = os.path.join(temp_dir, filename)
    chunks = extract_video(file_path=video_path, verbose=verbose, text_only=text_only)
    return chunks

def extract_audio(file_path: str, verbose: bool = False) -> List[Chunk]:
    import whisper # import only if needed
    model = whisper.load_model("base")
    result = model.transcribe(file_path, verbose=verbose)
    transcription = result['text']
    return [Chunk(path=file_path, text=transcription, image=None, source_type=SourceTypes.AUDIO)]

def extract_github(github_url: str, file_path: str = '', match: Optional[str] = None, ignore: Optional[str] = None, text_only: bool = False, ai_extraction: bool = False, branch: str = 'main', verbose: bool = False) -> List[Chunk]:
    files_contents = []
    if not GITHUB_TOKEN:
        raise ValueError("GITHUB_TOKEN environment variable is not set.")
    # make new tempdir for cloned repo
    with tempfile.TemporaryDirectory() as temp_dir:
        os.system(f"git clone {github_url} {temp_dir} --quiet")
        files_contents = extract_from_directory(dir_path=temp_dir, match=match, ignore=ignore, verbose=verbose, ai_extraction=ai_extraction, text_only=text_only)
    return files_contents

def extract_docx(file_path: str, verbose: bool = False, text_only: bool = False) -> List[Chunk]:
    # make new temp image directory
    import docx2txt # import only if needed
    chunks = []
    temp_image_dir = tempfile.mkdtemp()
    text = docx2txt.process(file_path, temp_image_dir)
    chunks.append(Chunk(path=file_path, text=text, image=None, source_type=SourceTypes.DOCX))
    if not text_only:
        for image_name in os.listdir(temp_image_dir):
            image_path = os.path.join(temp_image_dir, image_name)
            image = Image.open(image_path)
            image.load() # needed to close the file
            chunks.append(Chunk(path=file_path, text=None, image=image, source_type=SourceTypes.DOCX))
        # if temp dir exists, remove images and it
        if os.path.exists(temp_image_dir):
            for image_name in os.listdir(temp_image_dir):
                image_path = os.path.join(temp_image_dir, image_name)
                os.remove(image_path)
            os.rmdir(temp_image_dir)
    return chunks

def extract_pptx(file_path: str, verbose: bool = False, text_only: bool = False) -> List[Chunk]:
    from pptx import Presentation # import only if needed
    from pptx.enum.shapes import MSO_SHAPE_TYPE # import only if needed
    prs = Presentation(file_path)
    chunks = []
    # parse shapes inside slides
    for slide in prs.slides:
        slide_text = ""
        slide_images = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_text += f"{paragraph.text}\n\n"
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and not text_only:
                image_data = shape.image.blob
                image = Image.open(BytesIO(image_data))
                slide_images.append(image)
        chunks.append(Chunk(path=file_path, text=slide_text, image=None, source_type=SourceTypes.PPTX))
        for image in slide_images:
            chunks.append(Chunk(path=file_path, text=None, image=image, source_type=SourceTypes.PPTX))
    return chunks

def extract_from_ipynb(file_path: str, verbose: bool = False, ai_extraction: bool = False, text_only: bool = False) -> List[Chunk]:
    with open(file_path, 'r', encoding='utf-8') as file:
        notebook = json.load(file)
    chunks = []
    for cell in notebook['cells']:
        if cell['cell_type'] == 'markdown':
            text = ''.join(cell['source'])
            if not text_only:
                text, image_chunks = extract_images_from_markdown(text)
                chunks += image_chunks
            chunks.append(Chunk(path=file_path, text=text, image=None, source_type=SourceTypes.IPYNB))
        elif cell['cell_type'] == 'code':
            source = ''.join(cell['source'])
            output_text = ''
            if 'outputs' in cell:
                for output in cell['outputs']:
                    if 'data' in output and 'image/png' in output['data'] and not text_only:
                        image_data = output['data']['image/png']
                        image = Image.open(BytesIO(base64.b64decode(image_data)))
                        chunks.append(Chunk(path=file_path, text=None, image=image, source_type=SourceTypes.IPYNB))
                    elif 'data' in output and 'text/plain' in output['data']:
                        output_text += ''.join(output['data']['text/plain'])
            formatted_source = f"Source:\n```python\n{source}\n```\n"
            formatted_output = f"Output:\n```\n{output_text}\n```\n"
            chunks.append(Chunk(path=file_path, text=formatted_source+formatted_output, image=None, source_type=SourceTypes.IPYNB))
        elif cell['cell_type'] == 'raw':
            text = ''.join(cell['source'])
            chunks.append(Chunk(path=file_path, text=text, image=None, source_type=SourceTypes.IPYNB))
    return chunks