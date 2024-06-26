import base64
from concurrent.futures import ThreadPoolExecutor
from io import BytesIO
import math
import re
from typing import Dict, List, Optional, Tuple
import glob
import os
import tempfile
from urllib.parse import urlparse
import zipfile
from PIL import Image
import requests
import json
from .core import Chunk
import tempfile
import mimetypes
import dotenv
from datetime import datetime
from json import JSONEncoder
import csv
import shutil
from magika import Magika
dotenv.load_dotenv()

class JSONDateEncoder(JSONEncoder):
    def default(self, obj):
        if isinstance(obj, datetime):
            return obj.isoformat()
        return JSONEncoder.default(self, obj)

FOLDERS_TO_IGNORE = ['*node_modules.*', '.*venv.*', '.*.git.*', '.*.vscode.*', '.*pycache.*', '.*dist.*', '.*build.*', '.*target.*', '.*out.*', '.*output.*', '.*outputs*']
FILES_TO_IGNORE = ['package-lock.json', '.gitignore', '.*.bin', '.*.pyc', '.*.pyo', '.*.exe', '.*.bat', '.*.dll', '.*.obj', '.*.o', '.*.a', '.*.lib', '.*.so', '.*.dylib', '.*.ncb', '.*.sdf', '.*.suo', '.*.pdb', '.*.idb', '.*.pyd', '.*.ipynb_checkpoints', '.*.npy', '.*.pth']
GITHUB_TOKEN: str = os.getenv("GITHUB_TOKEN", None)
THEPIPE_API_KEY: str = os.getenv("THEPIPE_API_KEY", None)
USER_AGENT_STRING: str = os.getenv("USER_AGENT_STRING", "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/58.0.3029.110 Safari/537.3")
MAX_WHISPER_DURATION = 1200 # 20 minutes
TWITTER_DOMAINS = ['https://twitter.com', 'https://www.twitter.com', 'https://x.com', 'https://www.x.com']
YOUTUBE_DOMAINS = ['https://www.youtube.com', 'https://youtube.com']
GITHUB_DOMAINS = ['https://github.com', 'https://www.github.com']

def detect_source_type(source: str) -> str:
    # otherwise, try to detect the file type by its extension
    _, extension = os.path.splitext(source)
    if extension:
        guessed_mimetype = mimetypes.guess_type(source)[0]
        if guessed_mimetype:
            return guessed_mimetype
    # if that fails, try AI detection with Magika
    magika = Magika()
    with open(source, 'rb') as file:
        result = magika.identify_bytes(file.read())
    mimetype = result.output.mime_type
    return mimetype

def scrape_file(source: str, verbose: bool = False, ai_extraction: bool = False, text_only: bool = False) -> List[Chunk]:
    # returns chunks of scraped content from any source (file, URL, etc.)
    extraction = []
    source_type = detect_source_type(source)
    if source_type is None:
        if verbose:
            print(f"[thepipe] Unsupported source type: {source}", status="error")
        return extraction
    if verbose: 
        print(f"[thepipe] Scraping {source_type}: {source}...", status="info")
    if source_type == 'application/pdf':
        extraction = scrape_pdf(file_path=source, ai_extraction=ai_extraction, text_only=text_only, verbose=verbose)
    elif source_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
        extraction = scrape_docx(file_path=source, verbose=verbose, text_only=text_only)
    elif source_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
        extraction = scrape_pptx(file_path=source, verbose=verbose, text_only=text_only)
    elif source_type.startswith('image/'):
        extraction = scrape_image(file_path=source, text_only=text_only)
    elif source_type.startswith('application/vnd.ms-excel') or source_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
        extraction = scrape_spreadsheet(file_path=source)
    elif source_type == 'application/x-ipynb+json':
        extraction = scrape_ipynb(file_path=source, verbose=verbose, text_only=text_only)
    elif source_type == 'application/zip':
        extraction = scrape_zip(file_path=source, verbose=verbose, ai_extraction=ai_extraction, text_only=text_only)
    elif source_type.startswith('video/'):
        extraction = scrape_video(file_path=source, verbose=verbose, text_only=text_only)
    elif source_type.startswith('audio/'):
        extraction = scrape_audio(file_path=source, verbose=verbose)
    elif source_type.startswith('text/'):
        extraction = scrape_plaintext(file_path=source)
    else:
        try:
            extraction = scrape_plaintext(file_path=source)
        except Exception as e:
            if verbose: 
                print(f"[thepipe] Error extracting from {source}: {e}")
    if verbose: 
        if extraction:
            print(f"[thepipe] Extracted from {source}", status="success")
        else:
            print(f"[thepipe] No content extracted from {source}", status="error")
    
    return extraction

def scrape_plaintext(file_path: str) -> List[Chunk]:
    with open(file_path, 'r', encoding='utf-8') as file:
        text = file.read()
    return [Chunk(path=file_path, texts=[text])]

def scrape_directory(dir_path: str, include_regex: Optional[List[str]] = [], ignore_regex: Optional[List[str]] = [], verbose: bool = False, ai_extraction: bool = False, text_only: bool = False) -> List[Chunk]:
    extraction = []
    all_files = glob.glob(dir_path + "/**/*", recursive=True)
    for include_regex_string in include_regex:
        all_files = [file for file in all_files if re.search(include_regex_string, file, re.IGNORECASE)]
    for ignore_regex_string in ignore_regex + FILES_TO_IGNORE + FOLDERS_TO_IGNORE:
        all_files = [file for file in all_files if not re.search(ignore_regex_string, file, re.IGNORECASE)]
    with ThreadPoolExecutor() as executor:
        results = executor.map(lambda file_path: scrape_file(source=file_path, match=include_regex, ignore=ignore_regex, verbose=verbose, ai_extraction=ai_extraction, text_only=text_only), all_files)
        for result in results:
            extraction += result
    return extraction

def scrape_zip(file_path: str, match: Optional[str] = None, ignore: Optional[str] = None, verbose: bool = False, ai_extraction: bool = False, text_only: bool = False) -> List[Chunk]:
    extracted_files = []
    with tempfile.TemporaryDirectory() as temp_dir:
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        extracted_files = scrape_directory(dir_path=temp_dir, include_regex=match, ignore_regex=ignore, verbose=verbose, ai_extraction=ai_extraction, text_only=text_only)
    return extracted_files

def scrape_pdf(file_path: str, ai_extraction: bool = False, text_only: bool = False, verbose: bool = False) -> List[Chunk]:
    chunks = []

    if ai_extraction:
        # ai_extraction uses layout analysis AI to extract markdown, equations, tables, and images from the PDF
        MD_FOLDER = 'mdoutput'
        if not os.path.exists(MD_FOLDER):
            os.makedirs(MD_FOLDER)
        else:
            shutil.rmtree(MD_FOLDER)
            os.makedirs(MD_FOLDER)
        os.system(f"marker_single {file_path} {MD_FOLDER} --batch_multiplier 4 --max_pages 1000 --langs English")
        # Find the .md file and read its content
        for output_file in glob.glob(f'{MD_FOLDER}/*/*', recursive=True):
            if output_file.endswith('.md'):
                with open(output_file, 'r') as f:
                    markdown = f.read()
                    break
        if text_only:
            chunks.append(Chunk(path=file_path, texts=[markdown]))
            return chunks
        # split the markdown into text and images, so we can return them in the correct order
        content_pattern = re.compile(r'(\!\[.*?\]\(.*?\)|[^!\[]+)')
        content_matches = content_pattern.findall(markdown)
        for content in content_matches:
            if content.startswith('!['):
                # matched an image
                if text_only:
                    continue
                image_url = re.search(r'\((.*?)\)', content).group(1)
                try:
                    image = Image.open(requests.get(image_url, stream=True).raw)
                    chunks.append(Chunk(path=file_path, images=[image]))
                except Exception as e:
                    if verbose: print(f"[thepipe] Error loading image {image_url}: {e}", status='error')
            else:
                # matched text
                chunks.append(Chunk(path=file_path, texts=[content.strip()]))
        # remove the output folder
        shutil.rmtree(MD_FOLDER)
        if verbose: print(f"[thepipe] AI extracted from {file_path}", status='success')
        return chunks
    else:
        # if not using AI extraction, for each page, extract markdown and (optionally) full page images
        import fitz
        import pymupdf4llm
        doc = fitz.open(file_path)
        md_reader = pymupdf4llm.helpers.pymupdf_rag.to_markdown(doc, page_chunks=True)
        for i, page in enumerate(doc):
            #text = page.get_text()
            text = md_reader[i]["text"]
            if text_only:
                chunks.append(Chunk(path=file_path, texts=[text]))
            else:
                pix = page.get_pixmap()
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                chunks.append(Chunk(path=file_path, texts=[text], images=[img]))
        doc.close()
    return chunks

def get_images_from_ipynb_markdown(text: str) -> List[Image.Image]:
    image_urls = re.findall(r"!\[.*?\]\((.*?)\)", text)
    images = []
    for url in image_urls:
        extension = os.path.splitext(urlparse(url).path)[1]
        if extension in {'.jpg', '.jpeg', '.png'}:
            img = Image.open(requests.get(url, stream=True).raw)
        else:
            # ignore incompatible image extractions
            continue
        images.append(img)
    return images

def scrape_image(file_path: str, text_only: bool = False) -> List[Chunk]:
    import pytesseract
    img = Image.open(file_path)
    img.load()  # needed to close the file
    chunks = []
    if text_only:
        text = pytesseract.image_to_string(img)
        chunks.append(Chunk(path=file_path, texts=[text]))
    else:
        chunks.append(Chunk(path=file_path, images=[img]))
    return chunks

def scrape_spreadsheet(file_path: str) -> List[Chunk]:
    import pandas as pd
    if file_path.endswith(".csv"):
        df = pd.read_csv(file_path)
    elif file_path.endswith(".xls") or file_path.endswith(".xlsx"):
        df = pd.read_excel(file_path)
    else:
        raise ValueError("Unsupported file format")
    dicts = df.to_dict(orient='records')
    chunks = []
    for i, item in enumerate(dicts):
        # format each row as json along with the row index
        item['row index'] = i
        item_json = json.dumps(item, indent=4)
        chunks.append(Chunk(path=file_path, texts=[item_json]))
    return chunks

def extract_page_content(url: str, verbose: bool = False) -> Tuple[str, List[str]]:
    from urllib.parse import urlparse
    import markdownify
    from bs4 import BeautifulSoup
    from playwright.sync_api import sync_playwright
    with sync_playwright() as p:
        browser = p.chromium.launch()
        context = browser.new_context(user_agent=USER_AGENT_STRING)
        page = context.new_page()
        page.goto(url, wait_until='domcontentloaded')
        # Scroll to the bottom of the page to load dynamic content
        viewport_height = page.viewport_size['height']
        total_height = page.evaluate("document.body.scrollHeight")
        current_scroll_position = 0
        scrolldowns, max_scrolldowns = 0, 20  # Finite to prevent infinite scroll
        while current_scroll_position < total_height and scrolldowns < max_scrolldowns:
            page.wait_for_timeout(100)  # Wait for dynamic content to load
            current_scroll_position += viewport_height
            page.evaluate(f"window.scrollTo(0, {current_scroll_position})")
            scrolldowns += 1
            total_height = page.evaluate("document.body.scrollHeight")
        # Extract HTML content
        html_content = page.content()
        # Convert HTML to Markdown
        soup = BeautifulSoup(html_content, 'html.parser')
        markdown_content = markdownify.markdownify(str(soup), heading_style="ATX")
        # remove excessive newlines in the markdown
        while '\n\n\n' in markdown_content:
            markdown_content = markdown_content.replace('\n\n\n', '\n\n')
        # Extract images from the page
        images = []
        for img in page.query_selector_all('img'):
            img_path = img.get_attribute('src')
            if img_path and not img_path.startswith('data:image'):
                try:
                    if 'https' not in img_path:
                        base_url = urlparse(url).scheme + '://' + urlparse(url).netloc
                        img_path = base_url + img_path
                    images.append(img_path)
                except Exception as e:
                    if verbose: print(f"[thepipe] Ignoring error loading image {img_path}: {e}", status='error')
                    continue # ignore incompatible image extractions
        browser.close()

    return markdown_content, images

def parse_html_to_markdown(html_content):
    from bs4 import BeautifulSoup, NavigableString, Tag
    soup = BeautifulSoup(html_content, 'html.parser')
    markdown_content = []
    # recursively traverse the HTML content and extract text and links
    def traverse_and_extract(element):
        if isinstance(element, NavigableString):
            markdown_content.append(str(element))
        elif isinstance(element, Tag):
            if element.name == 'a' and 'href' in element.attrs:
                link_text = element.get_text()
                link_url = element['href']
                markdown_content.append(f'[{link_text}]({link_url})')
            else:
                for child in element.children:
                    traverse_and_extract(child)
    # extract content from the body tag
    body = soup.body
    if body:
        traverse_and_extract(body)
    return ''.join(markdown_content)

def scrape_url(url: str, text_only: bool = False, ai_extraction: bool = False) -> List[Chunk]:
    if any(url.startswith(domain) for domain in TWITTER_DOMAINS):
        extraction = scrape_tweet(url=url, text_only=text_only)
        return extraction
    elif any(url.startswith(domain) for domain in YOUTUBE_DOMAINS):
        extraction = scrape_youtube(youtube_url=url, text_only=text_only)
        return extraction
    elif any(url.startswith(domain) for domain in GITHUB_DOMAINS):
        extraction = scrape_github(github_url=url, text_only=text_only, ai_extraction=ai_extraction)
        return extraction
    _, extension = os.path.splitext(urlparse(url).path)
    all_texts = []
    all_images = []
    if extension and extension not in {'.html', '.htm', '.php', '.asp', '.aspx'}:
        # if url leads to a file, attempt to download it and scrape it
        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = os.path.join(temp_dir, os.path.basename(url))
            response = requests.get(url)
            with open(file_path, 'wb') as file:
                file.write(response.content)
            chunks = scrape_file(source=file_path, ai_extraction=ai_extraction, text_only=text_only)
            for chunk in chunks:
                all_texts.extend(chunk.texts)
                all_images.extend(chunk.images)
    else:
        # if url leads to web content, scrape it directly
        markdown_content, images = extract_page_content(url)
        all_texts.append(markdown_content)
        if not text_only:
            all_images.extend(images)
    if not all_texts and not all_images:
        raise ValueError("No content extracted from URL.")
    return [Chunk(path=url, texts=all_texts, images=all_images)]

def format_timestamp(seconds, chunk_index, chunk_duration):
    # helper function to format the timestamp.
    total_seconds = chunk_index * chunk_duration + seconds
    hours = int(total_seconds // 3600)
    minutes = int((total_seconds % 3600) // 60)
    seconds = total_seconds % 60
    milliseconds = int((seconds - int(seconds)) * 1000)
    return f"{hours:02}:{minutes:02}:{int(seconds):02}.{milliseconds:03}"

def scrape_video(file_path: str, verbose: bool = False, text_only: bool = False) -> List[Chunk]:
    import whisper
    from moviepy.editor import VideoFileClip

    model = whisper.load_model("base")
    video = VideoFileClip(file_path)
    num_chunks = math.ceil(video.duration / MAX_WHISPER_DURATION)
    chunks = []
    # split the video into chunks of fixed duration
    # here, we transcribe each chunk and extract its frame
    try:
        for i in range(num_chunks):
            start_time = i * MAX_WHISPER_DURATION
            end_time = start_time + MAX_WHISPER_DURATION
            if end_time > video.duration:
                end_time = video.duration
            # get the frame in the middle of the chunk
            frame_time = (start_time + end_time) / 2
            frame = video.get_frame(frame_time)
            image = Image.fromarray(frame)
            # save the audio to a temporary file
            with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as temp_audio_file:
                audio_path = temp_audio_file.name
            audio = video.subclip(start_time, end_time).audio
            transcription = None
            # transcribe it
            if audio is not None:
                audio.write_audiofile(audio_path, codec='pcm_s16le')
                result = model.transcribe(audio=audio_path, verbose=verbose)
                # Format transcription with timestamps
                formatted_transcription = []
                for segment in result['segments']:
                    start = format_timestamp(segment['start'], i, MAX_WHISPER_DURATION)
                    end = format_timestamp(segment['end'], i, MAX_WHISPER_DURATION)
                    formatted_transcription.append(f"[{start} --> {end}]  {segment['text']}")
                transcription = '\n'.join(formatted_transcription)
                os.remove(audio_path)
            texts = [transcription] if transcription else []
            images = [image] if not text_only else []
            if texts or images:
                chunks.append(Chunk(path=file_path, texts=texts, images=images))
    finally:
        video.close()
    return chunks

def scrape_youtube(youtube_url: str, text_only: bool = False, verbose: bool = False) -> List[Chunk]:
    from pytube import YouTube
    with tempfile.TemporaryDirectory() as temp_dir:
        filename = "temp_video.mp4"
        yt = YouTube(youtube_url)
        stream = yt.streams.filter(progressive=True, file_extension='mp4').first()
        stream.download(temp_dir, filename=filename)
        video_path = os.path.join(temp_dir, filename)
        # check if within max file size
        if os.path.getsize(video_path) > 10**8: # 100 MB
            raise ValueError("Video file is too large to process.")
        chunks = scrape_video(file_path=video_path, verbose=verbose, text_only=text_only)
    return chunks

def scrape_audio(file_path: str, verbose: bool = False) -> List[Chunk]:
    import whisper
    model = whisper.load_model("base")
    result = model.transcribe(audio=file_path, verbose=verbose)
    # Format transcription with timestamps
    formatted_transcription = []
    for segment in result['segments']:
        start = format_timestamp(segment['start'], 0, 0)
        end = format_timestamp(segment['end'], 0, 0)
        formatted_transcription.append(f"[{start} --> {end}]  {segment['text']}")
    # join the formatted transcription into a single string
    transcription = '\n'.join(formatted_transcription)
    return [Chunk(path=file_path, texts=[transcription])]

def scrape_github(github_url: str, include_regex: Optional[List[str]] = [], ignore_regex: Optional[List[str]] = [], text_only: bool = False, ai_extraction: bool = False, branch: str = 'main', verbose: bool = False) -> List[Chunk]:
    files_contents = []
    if not GITHUB_TOKEN:
        raise ValueError("GITHUB_TOKEN environment variable is not set.")
    # make new tempdir for cloned repo
    with tempfile.TemporaryDirectory() as temp_dir:
        # requires git
        os.system(f"git clone {github_url} {temp_dir} --quiet")
        files_contents = scrape_directory(dir_path=temp_dir, include_regex=include_regex, ignore_regex=ignore_regex, verbose=verbose, ai_extraction=ai_extraction, text_only=text_only)
    return files_contents

def scrape_docx(file_path: str, verbose: bool = False, text_only: bool = False) -> List[Chunk]:
    from docx import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table, _Cell
    from docx.text.paragraph import Paragraph
    import csv
    import io

    # helper function to iterate through blocks in the document
    def iter_block_items(parent):
        if parent.__class__.__name__ == 'Document':
            parent_elm = parent.element.body
        elif parent.__class__.__name__ == '_Cell':
            parent_elm = parent._tc
        else:
            raise ValueError("Unsupported parent type")
        # iterate through each child element in the parent element
        for child in parent_elm.iterchildren():
            if child.__class__.__name__ == 'CT_P':
                yield Paragraph(child, parent)
            elif child.__class__.__name__ == 'CT_Tbl':
                yield Table(child, parent)

    # helper function to read tables in the document
    def read_docx_tables(tab):
        vf = io.StringIO()
        writer = csv.writer(vf)
        for row in tab.rows:
            writer.writerow(cell.text for cell in row.cells)
        vf.seek(0)
        return vf.getvalue()

    # read the document
    document = Document(file_path)
    chunks = []
    image_counter = 0

    # Define namespaces
    nsmap = {
        'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
        'pic': 'http://schemas.openxmlformats.org/drawingml/2006/picture',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    }

    try:
        # scrape each block in the document to create chunks
        # A block can be a paragraph, table, or image
        for block in iter_block_items(document):
            block_texts = []
            block_images = []
            if block.__class__.__name__ == 'Paragraph':
                block_texts.append(block.text)
                if not text_only:
                    # "runs" are the smallest units in a paragraph
                    for run in block.runs:
                        if 'pic:pic' in run.element.xml:
                            # extract images from the paragraph
                            for pic in run.element.findall('.//pic:pic', nsmap):
                                cNvPr = pic.find('.//pic:cNvPr', nsmap)
                                name_attr = cNvPr.get("name") if cNvPr is not None else f"image_{image_counter}"
                                
                                blip = pic.find('.//a:blip', nsmap)
                                if blip is not None:
                                    embed_attr = blip.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
                                    if embed_attr:
                                        image_part = document.part.related_parts[embed_attr]
                                        image_data = io.BytesIO(image_part._blob)
                                        image = Image.open(image_data)
                                        block_images.append(image)
                                        image_counter += 1
            elif block.__class__.__name__ == 'Table':
                table_text = read_docx_tables(block)
                block_texts.append(table_text)
            if block_texts or block_images:
                chunks.append(Chunk(path=file_path, texts=block_texts, images=block_images))

    finally:
        # Close any open image files
        for chunk in chunks:
            for image in chunk.images:
                image.close()

    return chunks

def scrape_pptx(file_path: str, verbose: bool = False, text_only: bool = False) -> List[Chunk]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = Presentation(file_path)
    chunks = []
    # iterate through each slide in the presentation
    for slide in prs.slides:
        slide_texts = []
        slide_images = []
        # iterate through each shape in the slide
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text
                    if len(slide_texts) == 0:
                        text = '# ' + text # header for first text of a slide
                    slide_texts.append(text)
            # extract images from shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and not text_only:
                image_data = shape.image.blob
                image = Image.open(BytesIO(image_data))
                slide_images.append(image)
        # add slide to chunks if it has text or images
        if slide_texts or slide_images:
            chunks.append(Chunk(path=file_path, texts=slide_texts, images=slide_images))
    # return all chunks
    return chunks

def scrape_ipynb(file_path: str, verbose: bool = False, text_only: bool = False) -> List[Chunk]:
    with open(file_path, 'r', encoding='utf-8') as file:
        notebook = json.load(file)
    chunks = []
    # parse cells in the notebook
    for cell in notebook['cells']:
        texts = []
        images = []
        # parse cell content based on type
        if cell['cell_type'] == 'markdown':
            text = ''.join(cell['source'])
            if not text_only:
                images = get_images_from_ipynb_markdown(text)
            texts.append(text)
        elif cell['cell_type'] == 'code':
            source = ''.join(cell['source'])
            texts.append(source)
            output_texts = []
            # code cells can have outputs
            if 'outputs' in cell:
                for output in cell['outputs']:
                    if 'data' in output and 'image/png' in output['data'] and not text_only:
                        image_data = output['data']['image/png']
                        image = Image.open(BytesIO(base64.b64decode(image_data)))
                        images.append(image)
                    elif 'data' in output and 'text/plain' in output['data']:
                        output_text = ''.join(output['data']['text/plain'])
                        output_texts.append(output_text)
            if output_texts:
                texts.extend(output_texts)
        elif cell['cell_type'] == 'raw':
            text = ''.join(cell['source'])
            texts.append(text)
        if texts or images:
            chunks.append(Chunk(path=file_path, texts=texts, images=images))
    return chunks

def scrape_tweet(url: str, text_only: bool = False) -> List[Chunk]:
    # magic function from https://github.com/vercel/react-tweet/blob/main/packages/react-tweet/src/api/fetch-tweet.ts
    # unofficial, could break at any time
    def get_token(id: str) -> str:
        result = (float(id) / 1e15) * math.pi
        base_36_result = ''
        characters = '0123456789abcdefghijklmnopqrstuvwxyz'
        while result > 0:
            remainder = int(result % (6 ** 2))
            base_36_result = characters[remainder] + base_36_result
            result = (result - remainder) // (6 ** 2)
        base_36_result = re.sub(r'(0+|\.)', '', base_36_result)
        return base_36_result
    tweet_id = url.split('status/')[-1].split('?')[0]
    token = get_token(tweet_id)
    tweet_api_url = "https://cdn.syndication.twimg.com/tweet-result"
    params = {
        "id": tweet_id,
        "language": "en",
        "token": token
    }
    response = requests.get(tweet_api_url, params=params)
    if response.status_code != 200:
        raise ValueError(f"Failed to fetch tweet. Status code: {response.status_code}")
    tweet_data = response.json()
    # Extract tweet text
    tweet_text = tweet_data.get("text", "")
    # Extract images from tweet
    if not text_only:
        images = []
        if "mediaDetails" in tweet_data:
            for media in tweet_data["mediaDetails"]:
                image_url = media.get("media_url_https")
                if image_url:
                    image_response = requests.get(image_url)
                    img = Image.open(BytesIO(image_response.content))
                    images.append(img)
    # Create chunks for text and images
    chunk = Chunk(path=url, texts=[tweet_text], images=images)
    return [chunk]