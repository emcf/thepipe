import argparse
import base64
from enum import Enum
from io import BytesIO
import re
import time
from typing import *
import glob
import os
import tempfile
import unicodedata
from urllib.parse import unquote, urlparse
import zipfile
from colorama import Style, Fore
import pandas as pd
from PIL import Image
import requests
import json
import pytesseract
from unstructured.partition.auto import partition
from playwright.sync_api import sync_playwright
import fitz
from core import Chunk, print_status, SourceTypes
import docx2txt
import tempfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

FILES_TO_IGNORE = {'.gitignore', '.bin', '.pyc', '.pyo', '.exe', '.dll', '.obj', '.o', '.a', '.lib', '.so', '.dylib', '.ncb', '.sdf', '.suo', '.pdb', '.idb', '.pyd', '.ipynb_checkpoints', '.npy', '.pth'} # Files to ignore, please feel free to customize!
CODE_EXTENSIONS = {'.h', '.json', '.js', '.jsx',  '.cs', '.java', '.html', '.css', '.ini', '.xml', '.yaml', '.xaml', '.sh'} # Plaintext files that should not be compressed with LLMLingua
CTAGS_CODE_EXTENSIONS = {'.c', '.cpp', '.py', '.ts', '.tsx'} # code files that work with ctags
PLAINTEXT_EXTENSIONS = {'.txt', '.md', '.rtf'}
GITHUB_TOKEN: str = os.getenv("GITHUB_TOKEN")
SUPABASE_URL: str = os.getenv("SUPABASE_URL")
SUPABASE_KEY: str = os.getenv("SUPABASE_KEY")
MATHPIX_APP_ID: str = os.getenv("MATHPIX_APP_ID")
MATHPIX_APP_KEY: str = os.getenv("MATHPIX_APP_KEY")
MATHPIX_CALLS_PER_MIN: int = 6

def extract_from_source(source_string: str, match: Optional[str] = None, ignore: Optional[str] = None, limit: int = 64000, verbose: bool = False, mathpix: bool = False, text_only: bool = False) -> List[str]:
    source_type = detect_type(source_string)
    if source_type is None:
        return [Chunk(path=source_string)]
    if verbose: print_status(f"Extracting from {source_type.value}", status='info')
    if source_type == SourceTypes.DIR or source_string == '.' or source_string == './':
        if source_string == '.' or source_string == './':
            source_string = os.getcwd()
        return extract_from_directory(source_string=source_string, match=match, ignore=ignore, verbose=verbose, mathpix=mathpix, text_only=text_only)
    elif source_type == SourceTypes.GITHUB:
        return extract_github(github_url=source_string, file_path='', match=match, ignore=ignore, text_only=text_only, verbose=verbose, mathpix=mathpix, branch='master')
    elif source_type == SourceTypes.ZIP:
        return extract_zip(source_string=source_string, match=match, ignore=ignore, verbose=verbose, mathpix=mathpix, text_only=text_only)
    elif source_type == SourceTypes.URL:
        return [extract_url(url=source_string, text_only=text_only)]
    return extract_from_file(source_string=source_string, source_type=source_type, verbose=verbose, mathpix=mathpix, text_only=text_only)

def extract_from_file(source_string: str, source_type: str, verbose: bool = False, mathpix: bool = False, text_only: bool = False) -> List[str]:
    try:
        if source_type == SourceTypes.PDF:
            extraction = extract_pdf(source_string, mathpix, text_only)
        elif source_type == SourceTypes.DOCX:
            extraction = extract_docx(source_string)
        elif source_type == SourceTypes.PPTX:
            extraction = extract_pptx(source_string)
        elif source_type == SourceTypes.IMAGE:
            extraction = [extract_image(source_string, text_only)]
        elif source_type == SourceTypes.SPREADSHEET:
            extraction = [extract_spreadsheet(source_string)]
        elif source_type == SourceTypes.PLAINTEXT:
            extraction = [extract_plaintext(source_string)]
        elif source_type == SourceTypes.UNCOMPRESSIBLE_CODE:
            extraction = [extract_plaintext(source_string)]
            extraction = [Chunk(path=e.path, text=e.text, image=None, source_type=SourceTypes.UNCOMPRESSIBLE_CODE) for e in extraction] # change types to code
        elif source_type == SourceTypes.COMPRESSIBLE_CODE:
            extraction = [extract_plaintext(source_string)]
            extraction = [Chunk(path=e.path, text=e.text, image=None, source_type=SourceTypes.COMPRESSIBLE_CODE) for e in extraction]
        else:
            extraction = [extract_unstructured(source_string)]
        if verbose: print_status(f"Extracted from {source_string}", status='success')
        return extraction
    except Exception as e:
        if verbose: print_status(f"Failed to extract from {source_string}: {e}", status='error')
        return [Chunk(path=source_string)]

def detect_type(source_string: str) -> Optional[SourceTypes]:
    if source_string.startswith("https://github.com"):
        return SourceTypes.GITHUB
    elif source_string.startswith("http") or source_string.startswith("www.") or source_string.startswith("ftp."):
        return SourceTypes.URL
    elif source_string.endswith(".zip"):
        return SourceTypes.ZIP
    elif any(source_string.endswith(ext) for ext in CODE_EXTENSIONS):
        return SourceTypes.UNCOMPRESSIBLE_CODE
    elif any(source_string.endswith(ext) for ext in CTAGS_CODE_EXTENSIONS):
        return SourceTypes.COMPRESSIBLE_CODE
    elif source_string.endswith(".pdf"):
        return SourceTypes.PDF
    elif any(source_string.endswith(ext) for ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp', '.svg']):
        return SourceTypes.IMAGE
    elif any(source_string.endswith(ext) for ext in ['.csv', '.xls', '.xlsx']):
        return SourceTypes.SPREADSHEET
    elif source_string.endswith(".ipynb"):
        return SourceTypes.IPYNB
    elif source_string.endswith(".docx"):
        return SourceTypes.DOCX
    elif source_string.endswith(".pptx"):
        return SourceTypes.PPTX
    elif os.path.isdir(source_string) or source_string == '.' or source_string == './':
        return SourceTypes.DIR
    elif any(source_string.endswith(ext) for ext in PLAINTEXT_EXTENSIONS):
        return SourceTypes.PLAINTEXT
    else:
        return None # want to avoid processing unknown file types

def extract_unstructured(source_name: str) -> List[Chunk]:
    elements = partition(source_name)
    text = "\n\n".join([str(el) for el in elements])
    return Chunk(path=source_name, text=text, image=None, source_type=SourceTypes.PLAINTEXT)

def extract_plaintext(source_name: str) -> List[Chunk]:
    with open(source_name, 'r', encoding='utf-8') as file:
        text = file.read()
    return Chunk(path=source_name, text=text, image=None, source_type=SourceTypes.PLAINTEXT)

def extract_from_directory(source_string: str, match: Optional[str] = None, ignore: Optional[str] = None, verbose: bool = False, mathpix: bool = False, text_only: bool = False) -> List[Chunk]:
    all_files = glob.glob(source_string + "/**/*", recursive=True)
    matched_files = [file for file in all_files if re.search(match, file, re.IGNORECASE)] if match else all_files
    files_to_ignore = {file for file in matched_files if re.search(ignore, file, re.IGNORECASE)} if ignore else []
    file_paths = [file for file in matched_files if os.path.isfile(file) and file not in files_to_ignore]
    contents = []
    for file_path in file_paths:
        # ignore hidden files, modules, etc.
        if file_path.startswith('.') or any(file_path.endswith(extension) for extension in FILES_TO_IGNORE):
            continue
        # extract contents from the file
        contents += extract_from_source(source_string=file_path, match=match, ignore=ignore, verbose=verbose, mathpix=mathpix, text_only=text_only)
    return contents

def extract_zip(source_string: str, match: Optional[str] = None, ignore: Optional[str] = None, verbose: bool = False, mathpix: bool = False, text_only: bool = False) -> List[Chunk]:
    extracted_files = []
    with tempfile.TemporaryDirectory() as temp_dir:
        with zipfile.ZipFile(source_string, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        extracted_files = extract_from_directory(source_string=temp_dir, match=match, ignore=ignore, verbose=verbose, mathpix=mathpix, text_only=text_only)
    return extracted_files

def extract_pdf(source_name: str, mathpix: bool = False, text_only: bool = False) -> List[Chunk]:
    content_per_page = []
    if mathpix:
        headers = {
            "app_id": MATHPIX_APP_ID,
            "app_key": MATHPIX_APP_KEY,
            "Content-type": "application/json",
        }
        r = requests.post(
            "https://api.mathpix.com/v3/pdf",
            json={
                "url": source_name,
                "math_inline_delimiters": ["$", "$"],
                "math_display_delimiters": ["$$", "$$"]
            },
            headers=headers,
        )
        mathpix_id = r.json()["pdf_id"]
        processed_url = f"https://api.mathpix.com/v3/pdf/{mathpix_id}.mmd"
        response = requests.get(processed_url, headers=headers)
        for tries in range(10):
            time.sleep(60/MATHPIX_CALLS_PER_MIN)
            response = requests.get(processed_url, headers=headers)
            if response.status_code == 200:
                break
        if response.status_code != 200:
            raise Exception(f"Mathpix failed to process the PDF: {response.text}")
        # clean result to unicode error (normalize text and remove all special characters)
        text = response.text
        text = (
            unicodedata.normalize("NFKD", text)
            .encode("ASCII", "ignore")
            .decode("utf-8", "ignore")
        )
        content_per_page.append((text, None))
        # extract markdown images from the string
        image_urls = re.findall(r"!\[.*?\]\((.*?)\)", text)
        for url in image_urls:
            img = Image.open(requests.get(url, stream=True).raw)
            content_per_page.append(Chunk(path=url, text=None, image=img, source_type=SourceTypes.IMAGE))
    else:
        with open(source_name, 'rb') as file:
            doc = fitz.open(source_name)
            for page in doc:
                text = page.get_text()
                if not text_only:
                    pix = page.get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    content_per_page.append(Chunk(path=source_name, text=text, image=img,source_type=SourceTypes.PDF))
                else:
                    content_per_page.append(Chunk(path=source_name, text=text, image=None, source_type=SourceTypes.PDF))
            doc.close()
    return content_per_page

def extract_image(source_name: str, text_only: bool = False) -> List[Chunk]:
    img = Image.open(source_name)
    img.load() # needed to close the file
    if text_only:
        text = pytesseract.image_to_string(img)
        return Chunk(path=source_name, text=text, image=None, source_type=SourceTypes.IMAGE)
    else:
        return Chunk(path=source_name, text=None, image=img, source_type=SourceTypes.IMAGE)
    
def extract_spreadsheet(source_name: str) -> List[Chunk]:
    if source_name.endswith(".csv"):
        df = pd.read_csv(source_name)
    elif source_name.endswith(".xls") or source_name.endswith(".xlsx"):
        df = pd.read_excel(source_name)
    dict = df.to_dict(orient='records')
    json_dict = json.dumps(dict, indent=4)
    return Chunk(path=source_name, text=json_dict, image=None, source_type=SourceTypes.SPREADSHEET)
    
def extract_url(url: str, text_only: bool = False) -> List[Chunk]:
    img = None
    text = None
    with sync_playwright() as p:
        for browser_type in [p.chromium, p.firefox, p.webkit]:
            browser = browser_type.launch()
            page = browser.new_page()
            page.goto('https://scrapingant.com/')
            screenshot = page.screenshot()
            img = Image.open(BytesIO(screenshot))
            text = page.inner_text('body')
            browser.close()
    if img is None and text is None:
        raise Exception("Failed to extract from URL.")
    return Chunk(path=url, text=text, image=img, source_type=SourceTypes.URL)

def extract_github(github_url: str, file_path: str = '', match: Optional[str] = None, ignore: Optional[str] = None, text_only: bool = False, mathpix: bool = False, branch: str = 'main', verbose: bool = False) -> List[Chunk]:
    files_contents = []
    if not GITHUB_TOKEN:
        raise ValueError("GITHUB_TOKEN environment variable is not set.")
    # Function to extract repo details
    def extract_repo_details(url):
        path = urlparse(url).path
        path_parts = path.strip('/').split('/')
        if len(path_parts) >= 2:
            return path_parts[0], path_parts[1]
        else:
            raise ValueError("Invalid GitHub URL provided.")
    # Download response from GitHub API
    owner, repo = extract_repo_details(github_url)
    api_url = f"https://api.github.com/repos/{owner}/{repo}/contents/{unquote(file_path)}?ref={branch}"
    headers = {
        'Accept': 'application/vnd.github.v3.raw',
        'Authorization': f'token {GITHUB_TOKEN}',
        'X-GitHub-Api-Version': '2022-11-28'
    }
    response = requests.get(api_url, headers=headers)
    response_json = json.loads(response.text)
    # Navigate the response JSON
    for item in response_json:
        if 'path' not in item:
            continue
        path = item['path']
        if path.startswith('.') or any(path.endswith(extension) for extension in FILES_TO_IGNORE):
            continue
        if ignore and re.search(ignore, path, re.IGNORECASE):
            continue
        if item['type'] == 'file':
            # get the file content
            file_content_request = requests.get(item['download_url'], headers=headers)
            extension = os.path.splitext(path)[1].lower()
            temp_file_path = f"temp{extension}"
            # write the file content to a temporary file
            with open(temp_file_path, 'wb') as temp_file:
                temp_file.write(file_content_request.content)
                source_type = detect_type(temp_file_path)
                if source_type is not None:
                    extractions = extract_from_file(source_string=temp_file_path, source_type=source_type, text_only=text_only, verbose=verbose, mathpix=mathpix)
                    for extraction in extractions:
                        extraction.path = path # use github path, not temp path
                    files_contents += extractions
            os.remove(temp_file_path)
        elif item['type'] == 'dir':
            files_contents += extract_github(github_url=github_url, file_path=path, match=match, text_only=text_only, mathpix=mathpix, branch=branch, verbose=verbose)
    return files_contents

def extract_docx(source_name: str) -> List[Chunk]:
    # make new temp image directory
    chunks = []
    temp_image_dir = tempfile.mkdtemp()
    print('processing')
    text = docx2txt.process(source_name, temp_image_dir)
    chunks.append(Chunk(path=source_name, text=text, image=None, source_type=SourceTypes.DOCX))
    for image_name in os.listdir(temp_image_dir):
        print(image_name)
        image_path = os.path.join(temp_image_dir, image_name)
        print('attempgin to open')
        image = Image.open(image_path)
        image.load() # needed to close the file
        print("appending")
        chunks.append(Chunk(path=source_name, text=None, image=image, source_type=SourceTypes.DOCX))
    # if temp dir exists, remove images and it
    print('attempting delete')
    if os.path.exists(temp_image_dir):
        for image_name in os.listdir(temp_image_dir):
            image_path = os.path.join(temp_image_dir, image_name)
            os.remove(image_path)
        os.rmdir(temp_image_dir)
    print('done')
    return chunks

def extract_pptx(source_name: str) -> List[Chunk]:
    prs = Presentation(source_name)
    chunks = []
    # parse slides, shapes, and images
    for slide in prs.slides:
        slide_text = ""
        slide_images = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_text += f"{paragraph.text}\n\n"
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_data = shape.image.blob
                image = Image.open(BytesIO(image_data))
                slide_images.append(image)
        chunks.append(Chunk(path=source_name, text=slide_text, image=None, source_type=SourceTypes.PPTX))
        for image in slide_images:
            chunks.append(Chunk(path=source_name, text=None, image=image, source_type=SourceTypes.PPTX))
    return chunks