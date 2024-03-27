import argparse
import base64
from concurrent.futures import ThreadPoolExecutor
from enum import Enum
from io import BytesIO
import re
import time
from typing import *
import glob
import os
import tempfile
import unicodedata
from urllib.parse import unquote, urlparse
import zipfile
from colorama import Style, Fore
import pandas as pd
from PIL import Image
import requests
import json
import pytesseract
from unstructured.partition.auto import partition
from playwright.sync_api import sync_playwright
import fitz
from core import Chunk, print_status, SourceTypes
import docx2txt
import tempfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

FILES_TO_IGNORE = {'package-lock.json', '.gitignore', '.bin', '.pyc', '.pyo', '.exe', '.dll', '.obj', '.o', '.a', '.lib', '.so', '.dylib', '.ncb', '.sdf', '.suo', '.pdb', '.idb', '.pyd', '.ipynb_checkpoints', '.npy', '.pth'} # Files to ignore, please feel free to customize!
CODE_EXTENSIONS = {'.h', '.json', '.js', '.jsx', '.ts', '.tsx',  '.cs', '.java', '.html', '.css', '.ini', '.xml', '.yaml', '.xaml', '.sh'} # Plaintext files that should not be compressed with LLMLingua
CTAGS_CODE_EXTENSIONS = {'.c', '.cpp', '.py'} # code files that work with ctags
PLAINTEXT_EXTENSIONS = {'.txt', '.md', '.rtf'}
GITHUB_TOKEN: str = os.getenv("GITHUB_TOKEN")
SUPABASE_URL: str = os.getenv("SUPABASE_URL")
SUPABASE_KEY: str = os.getenv("SUPABASE_KEY")
MATHPIX_APP_ID: str = os.getenv("MATHPIX_APP_ID")
MATHPIX_APP_KEY: str = os.getenv("MATHPIX_APP_KEY")
MATHPIX_CALLS_PER_MIN: int = 6

def extract_from_source(source: str, match: Optional[str] = None, ignore: Optional[str] = None, limit: int = 64000, verbose: bool = False, mathpix: bool = False, text_only: bool = False) -> List[str]:
    source_type = detect_type(source)
    if source_type is None:
        return [Chunk(path=source)]
    if verbose: print_status(f"Extracting from {source_type.value}", status='info')
    if source_type == SourceTypes.DIR or source == '.' or source == './':
        if source == '.' or source == './':
            source = os.getcwd()
        return extract_from_directory(dir_path=source, match=match, ignore=ignore, verbose=verbose, mathpix=mathpix, text_only=text_only)
    elif source_type == SourceTypes.GITHUB:
        return extract_github(github_url=source, file_path='', match=match, ignore=ignore, text_only=text_only, verbose=verbose, mathpix=mathpix, branch='master')
    elif source_type == SourceTypes.ZIP:
        return extract_zip(file_path=source, match=match, ignore=ignore, verbose=verbose, mathpix=mathpix, text_only=text_only)
    elif source_type == SourceTypes.URL:
        return [extract_url(url=source, text_only=text_only)]
    elif source_type == SourceTypes.IPYNB:
        return extract_from_ipynb(file_path=source, verbose=verbose, mathpix=mathpix, text_only=text_only)
    return extract_from_file(file_path=source, source_type=source_type, verbose=verbose, mathpix=mathpix, text_only=text_only)

def extract_from_file(file_path: str, source_type: str, verbose: bool = False, mathpix: bool = False, text_only: bool = False) -> List[str]:
    try:
        if source_type == SourceTypes.PDF:
            extraction = extract_pdf(file_path=file_path, mathpix=mathpix, text_only=text_only, verbose=verbose)
        elif source_type == SourceTypes.DOCX:
            extraction = extract_docx(file_path=file_path, verbose=verbose, text_only=text_only)
        elif source_type == SourceTypes.PPTX:
            extraction = extract_pptx(file_path=file_path, verbose=verbose, text_only=text_only)
        elif source_type == SourceTypes.IMAGE:
            extraction = [extract_image(file_path=file_path, text_only=text_only)]
        elif source_type == SourceTypes.SPREADSHEET:
            extraction = [extract_spreadsheet(file_path=file_path)]
        elif source_type == SourceTypes.PLAINTEXT:
            extraction = [extract_plaintext(file_path=file_path)]
        elif source_type == SourceTypes.UNCOMPRESSIBLE_CODE:
            extraction = [extract_plaintext(file_path=file_path)]
            extraction = [Chunk(path=e.path, text=e.text, image=None, source_type=SourceTypes.UNCOMPRESSIBLE_CODE) for e in extraction] # change types to code
        elif source_type == SourceTypes.COMPRESSIBLE_CODE:
            extraction = [extract_plaintext(file_path=file_path)]
            extraction = [Chunk(path=e.path, text=e.text, image=None, source_type=SourceTypes.COMPRESSIBLE_CODE) for e in extraction]
        else:
            extraction = [extract_unstructured(file_path=file_path)]
        if verbose: print_status(f"Extracted from {file_path}", status='success')
        return extraction
    except Exception as e:
        if verbose: print_status(f"Failed to extract from {file_path}: {e}", status='error')
        return [Chunk(path=file_path)]

def detect_type(source: str) -> Optional[SourceTypes]:
    if source.startswith("https://github.com"):
        return SourceTypes.GITHUB
    elif source.startswith("http") or source.startswith("www.") or source.startswith("ftp."):
        return SourceTypes.URL
    elif source.endswith(".zip"):
        return SourceTypes.ZIP
    elif any(source.endswith(ext) for ext in CODE_EXTENSIONS):
        return SourceTypes.UNCOMPRESSIBLE_CODE
    elif any(source.endswith(ext) for ext in CTAGS_CODE_EXTENSIONS):
        return SourceTypes.COMPRESSIBLE_CODE
    elif source.endswith(".pdf"):
        return SourceTypes.PDF
    elif any(source.endswith(ext) for ext in ['.jpg', '.jpeg', '.png',]): # TODO: '.svg', '.webp', '.gif', '.bmp', '.tiff'
        return SourceTypes.IMAGE
    elif any(source.endswith(ext) for ext in ['.csv', '.xls', '.xlsx']):
        return SourceTypes.SPREADSHEET
    elif source.endswith(".ipynb"):
        return SourceTypes.IPYNB
    elif source.endswith(".docx"):
        return SourceTypes.DOCX
    elif source.endswith(".pptx"):
        return SourceTypes.PPTX
    elif os.path.isdir(source) or source == '.' or source == './':
        return SourceTypes.DIR
    elif any(source.endswith(ext) for ext in PLAINTEXT_EXTENSIONS):
        return SourceTypes.PLAINTEXT
    else:
        return None # want to avoid processing unknown file types

def extract_unstructured(file_path: str) -> List[Chunk]:
    elements = partition(file_path)
    text = "\n\n".join([str(el) for el in elements])
    return Chunk(path=file_path, text=text, image=None, source_type=SourceTypes.PLAINTEXT)

def extract_plaintext(file_path: str) -> List[Chunk]:
    with open(file_path, 'r', encoding='utf-8') as file:
        text = file.read()
    return Chunk(path=file_path, text=text, image=None, source_type=SourceTypes.PLAINTEXT)

def should_ignore(file_path: str, ignore: Optional[str] = None) -> bool:
    if ignore is not None and re.search(ignore, file_path, re.IGNORECASE):
        return True
    if any(file_path.endswith(extension) for extension in FILES_TO_IGNORE):
        return True
    if file_path.startswith('.'):
        return True
    if any(x in file_path for x in ['node_modules', 'venv', '.git', '.vscode', '__pycache__']):
        return True
    if not os.path.isfile(file_path):
        return True
    return False

def extract_from_directory(dir_path: str, match: Optional[str] = None, ignore: Optional[str] = None, verbose: bool = False, mathpix: bool = False, text_only: bool = False) -> List[Chunk]:
    all_files = glob.glob(dir_path + "/**/*", recursive=True)
    matched_files = [file for file in all_files if re.search(match, file, re.IGNORECASE)] if match else all_files
    file_paths = [file for file in matched_files if not should_ignore(file, ignore)]
    contents = []
    with ThreadPoolExecutor() as executor:
        results = executor.map(lambda file_path: extract_from_source(source=file_path, match=match, ignore=ignore, verbose=verbose, mathpix=mathpix, text_only=text_only), file_paths)
        for result in results:
            contents += result
    return contents

def extract_zip(file_path: str, match: Optional[str] = None, ignore: Optional[str] = None, verbose: bool = False, mathpix: bool = False, text_only: bool = False) -> List[Chunk]:
    extracted_files = []
    with tempfile.TemporaryDirectory() as temp_dir:
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        extracted_files = extract_from_directory(dir_path=temp_dir, match=match, ignore=ignore, verbose=verbose, mathpix=mathpix, text_only=text_only)
    return extracted_files

def extract_pdf(file_path: str, mathpix: bool = False, text_only: bool = False, verbose: bool = False) -> List[Chunk]:
    chunks = []
    if mathpix:
        base_url = "https://api.mathpix.com/v3/pdf/"
        # extract text images, equations, and tables from the PDF using Mathpix
        headers = {
            "app_id": MATHPIX_APP_ID,
            "app_key": MATHPIX_APP_KEY,
        }
        data = {"options_json": json.dumps({
            "conversion_formats": {"md": True},
        })}
        with open(file_path, "rb") as f:
            files = {"file": f}
            response = requests.post(base_url, headers=headers, files=files, data=data)
        response_data = response.json()
        pdf_id = response_data["pdf_id"]
        # check if the processing is completed every 5 seconds
        for _ in range(10):
            response = requests.get(base_url + pdf_id, headers=headers)
            response_data = response.json()
            status = response_data.get("status", None)
            if status == "completed":
                response = requests.get(f"{base_url}{pdf_id}.md", headers=headers)
                 # clean result to unicode error (normalize text and remove all special characters)
                text = response.content.decode("utf-8").encode("ASCII", "ignore").decode("utf-8", "ignore")
                chunks.append(Chunk(path=file_path, text=text, image=None, source_type=SourceTypes.PDF))
                # extract markdown images from the 
                if not text_only:
                    chunks += extract_images_from_markdown(text)
                return chunks
            elif status == "error":
                raise ValueError("Unable to retrieve PDF from Mathpix")
            else:
                if verbose: print_status(f"Waiting for processing to complete...")
                time.sleep(5)
        raise TimeoutError("Mathpix processing took too long.")
    else:
        # extract text and images of each page from the PDF
        with open(file_path, 'rb') as file:
            doc = fitz.open(file_path)
            for page in doc:
                text = page.get_text()
                if text_only:
                    chunks.append(Chunk(path=file_path, text=text, image=None, source_type=SourceTypes.PDF))
                else:
                    pix = page.get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    chunks.append(Chunk(path=file_path, text=text, image=img,source_type=SourceTypes.PDF))
            doc.close()
    return chunks

def extract_images_from_markdown(text: str) -> List[Chunk]:
    image_urls = re.findall(r"!\[.*?\]\((.*?)\)", text)
    images = []
    for url in image_urls:
        extension = os.path.splitext(urlparse(url).path)[1]
        if extension in ['.jpg', '.jpeg', '.png']:
            img = Image.open(requests.get(url, stream=True).raw)
        else:
            # ignore incompatible image extractions
            continue
        images.append(Chunk(path=url, text=None, image=img, source_type=SourceTypes.IMAGE))
    return images

def extract_image(file_path: str, text_only: bool = False) -> Chunk:
    img = Image.open(file_path)
    img.load() # needed to close the file
    if text_only:
        text = pytesseract.image_to_string(img)
        return Chunk(path=file_path, text=text, image=None, source_type=SourceTypes.IMAGE)
    else:
        return Chunk(path=file_path, text=None, image=img, source_type=SourceTypes.IMAGE)
    
def extract_spreadsheet(file_path: str) -> Chunk:
    if file_path.endswith(".csv"):
        df = pd.read_csv(file_path)
    elif file_path.endswith(".xls") or file_path.endswith(".xlsx"):
        df = pd.read_excel(file_path)
    dict = df.to_dict(orient='records')
    json_dict = json.dumps(dict, indent=4)
    return Chunk(path=file_path, text=json_dict, image=None, source_type=SourceTypes.SPREADSHEET)
    
def extract_url(url: str, text_only: bool = False) -> Chunk:
    #os.system("python3 -m playwright install")
    img = None
    text = None
    with sync_playwright() as p:
        for browser_type in [p.chromium, p.firefox, p.webkit]:
            browser = browser_type.launch()
            page = browser.new_page()
            page.goto(url)
            img = None
            if not text_only:
                screenshot = page.screenshot()
                img = Image.open(BytesIO(screenshot))
            text = page.inner_text('body')
            browser.close()
    if img is None and text is None:
        raise Exception("Failed to extract from URL.")
    return Chunk(path=url, text=text, image=img, source_type=SourceTypes.URL)

def extract_github(github_url: str, file_path: str = '', match: Optional[str] = None, ignore: Optional[str] = None, text_only: bool = False, mathpix: bool = False, branch: str = 'main', verbose: bool = False) -> List[Chunk]:
    files_contents = []
    if not GITHUB_TOKEN:
        raise ValueError("GITHUB_TOKEN environment variable is not set.")
    # make new tempdir for cloned repo
    with tempfile.TemporaryDirectory() as temp_dir:
        os.system(f"git clone {github_url} {temp_dir} --quiet")
        files_contents = extract_from_directory(dir_path=temp_dir, match=match, ignore=ignore, verbose=verbose, mathpix=mathpix, text_only=text_only)
    return files_contents

def extract_docx(file_path: str, verbose: bool = False, text_only: bool = False) -> List[Chunk]:
    # make new temp image directory
    chunks = []
    temp_image_dir = tempfile.mkdtemp()
    text = docx2txt.process(file_path, temp_image_dir)
    chunks.append(Chunk(path=file_path, text=text, image=None, source_type=SourceTypes.DOCX))
    if not text_only:
        for image_name in os.listdir(temp_image_dir):
            image_path = os.path.join(temp_image_dir, image_name)
            image = Image.open(image_path)
            image.load() # needed to close the file
            chunks.append(Chunk(path=file_path, text=None, image=image, source_type=SourceTypes.DOCX))
        # if temp dir exists, remove images and it
        if os.path.exists(temp_image_dir):
            for image_name in os.listdir(temp_image_dir):
                image_path = os.path.join(temp_image_dir, image_name)
                os.remove(image_path)
            os.rmdir(temp_image_dir)
    return chunks

def extract_pptx(file_path: str, verbose: bool = False, text_only: bool = False) -> List[Chunk]:
    prs = Presentation(file_path)
    chunks = []
    # parse shapes inside slides
    for slide in prs.slides:
        slide_text = ""
        slide_images = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_text += f"{paragraph.text}\n\n"
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and not text_only:
                image_data = shape.image.blob
                image = Image.open(BytesIO(image_data))
                slide_images.append(image)
        chunks.append(Chunk(path=file_path, text=slide_text, image=None, source_type=SourceTypes.PPTX))
        for image in slide_images:
            chunks.append(Chunk(path=file_path, text=None, image=image, source_type=SourceTypes.PPTX))
    return chunks

def extract_from_ipynb(file_path: str, verbose: bool = False, mathpix: bool = False, text_only: bool = False) -> List[Chunk]:
    with open(file_path, 'r', encoding='utf-8') as file:
        notebook = json.load(file)
    chunks = []
    for cell in notebook['cells']:
        if cell['cell_type'] == 'markdown':
            text = ''.join(cell['source'])
            if not text_only:
                chunks += extract_images_from_markdown(text)
            chunks.append(Chunk(path=file_path, text=text, image=None, source_type=SourceTypes.IPYNB))
        elif cell['cell_type'] == 'code':
            source = ''.join(cell['source'])
            output_text = ''
            if 'outputs' in cell:
                for output in cell['outputs']:
                    if 'data' in output and 'image/png' in output['data'] and not text_only:
                        image_data = output['data']['image/png']
                        image = Image.open(BytesIO(base64.b64decode(image_data)))
                        chunks.append(Chunk(path=file_path, text=None, image=image, source_type=SourceTypes.IPYNB))
                    elif 'data' in output and 'text/plain' in output['data']:
                        output_text += ''.join(output['data']['text/plain'])
            formatted_source = f"Source:\n```python\n{source}\n```\n"
            formatted_output = f"Output:\n```\n{output_text}\n```\n"
            chunks.append(Chunk(path=file_path, text=formatted_source+formatted_output, image=None, source_type=SourceTypes.IPYNB))
        elif cell['cell_type'] == 'raw':
            text = ''.join(cell['source'])
            chunks.append(Chunk(path=file_path, text=text, image=None, source_type=SourceTypes.IPYNB))
    return chunks