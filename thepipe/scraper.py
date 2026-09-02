from typing import Any, Callable, Dict, Iterable, List, Optional, Tuple, Union, cast
import base64
from concurrent.futures import ThreadPoolExecutor, as_completed
from collections import OrderedDict
from io import BytesIO, StringIO
import math
import re
import fnmatch
import os
import tempfile
from urllib.parse import urlparse
import zipfile
from PIL import Image
import requests
import json
from .core import (
    HOST_IMAGES,
    Chunk,
    make_image_url,
    DEFAULT_AI_MODEL,
)
from .chunker import (
    chunk_by_page,
    chunk_by_document,
    chunk_by_section,
    chunk_semantic,
    chunk_by_keywords,
    chunk_by_length,
    chunk_agentic,
)
import mimetypes
import dotenv
from magika import Magika
import markdownify
import fitz
from openai import OpenAI
from openai.types.chat.chat_completion_message_param import ChatCompletionMessageParam

dotenv.load_dotenv()

FOLDERS_TO_IGNORE = {
    "*node_modules*",
    "*.git*",
    "*venv*",
    "*.vscode*",
    "*pycache*",
    "*.ipynb_checkpoints",
}
FILES_TO_IGNORE = {
    ".gitignore",
    "*.bin",
    # Python compiled files
    "*.pyc",
    "*.pyo",
    "*.pyd",
    # Shared libraries and binaries
    "*.so",
    "*.dll",
    "*.exe",
    # Archives and packages
    "*.tar",
    "*.tar.gz",
    "*.egg-info",
    "package-lock.json",
    "package.json",
    # Lock, log, and metadata files
    "*.lock",
    "*.log",
    "Pipfile.lock",
    "requirements.lock",
    "*.exe",
    "*.dll",
    ".DS_Store",
    "Thumbs.db",
}
USER_AGENT_STRING: str = os.getenv(
    "USER_AGENT_STRING",
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/58.0.3029.110 Safari/537.3",
)
MAX_WHISPER_DURATION = int(os.getenv("MAX_WHISPER_DURATION", 600))  # 10 minutes

SCRAPING_PROMPT = os.getenv(
    "SCRAPING_PROMPT",
    """A document is given. Please output the entire extracted contents from the document in detailed markdown format.
Your accuracy is very important. Please be careful to not miss any content from the document.
Be sure to correctly output a comprehensive format markdown for all the document contents (including, but not limited to, headers, paragraphs, lists, tables, menus, equations, full text contents, titles, subtitles, appendices, page breaks, columns, footers, page numbers, watermarks, footnotes, captions, annotations, images, figures, charts, shapes, form fields, content controls, signatures, etc.)
Always reply immediately with only markdown.
Do not give the markdown in a code block. Simply output the raw markdown immediately.
Do not output anything else.""",
)


def _load_whisper():
    try:
        import whisper
    except ImportError as exc:  # pragma: no cover - optional dependency
        raise ImportError(
            "Audio and video transcription requires the optional dependency `openai-whisper`. "
            "Install it with `pip install thepipe-api[audio]` or include the `gpu` extra."
        ) from exc

    return whisper


def detect_source_mimetype(source: str) -> str:
    # try to detect the file type by its extension
    _, extension = os.path.splitext(source)
    if extension:
        if extension == ".ipynb":
            # special case for notebooks, mimetypes is not familiar
            return "application/x-ipynb+json"
        guessed_mimetype, _ = mimetypes.guess_type(source)
        if guessed_mimetype:
            return guessed_mimetype
    # if that fails, try AI detection with Magika
    magika = Magika()
    with open(source, "rb") as file:
        result = magika.identify_bytes(file.read())
    mimetype = result.output.mime_type
    return mimetype


def _is_within_directory(root_path: str, candidate_path: str) -> bool:
    try:
        return os.path.commonpath([root_path, candidate_path]) == root_path
    except ValueError:
        return False


def scrape_file(
    filepath: str,
    verbose: bool = False,
    chunking_method: Optional[Callable[[List[Chunk]], List[Chunk]]] = chunk_by_page,
    openai_client: Optional[OpenAI] = None,
    model: str = DEFAULT_AI_MODEL,
    include_input_images: bool = True,
    include_output_images: bool = True,
) -> List[Chunk]:
    """
    Scrapes a file and returns a list of Chunk objects containing the text and images extracted from the file.

    Parameters
    ----------
    filepath : str
        The path to the file to scrape.
    verbose : bool, optional
        If ``True``, prints verbose output.
    chunking_method : Callable, optional
        A function to chunk the scraped content. Defaults to chunk_by_page.
    openai_client : OpenAI, optional
        An OpenAI client instance for LLM processing. If provided, uses VLM to scrape PDFs.
    model : str, optional
        The LLM model name to use for processing. Defaults to DEFAULT_AI_MODEL.
    include_input_images : bool, optional
        If ``True``, includes input images in the messages sent to the LLM.
    include_output_images : bool, optional
        If ``True``, includes output images in the returned chunks.
    Returns
    -------
    List[Chunk]
        A list of Chunk objects containing the scraped content.
    """
    # returns chunks of scraped content from the given file
    scraped_chunks = []
    source_mimetype = detect_source_mimetype(filepath)
    if source_mimetype is None:
        if verbose:
            print(f"[thepipe] Unsupported source type: {filepath}")
        return scraped_chunks
    if verbose:
        print(f"[thepipe] Scraping {source_mimetype}: {filepath}...")
    if source_mimetype == "application/pdf":
        scraped_chunks = scrape_pdf(
            file_path=filepath,
            verbose=verbose,
            model=model,
            openai_client=openai_client,
            include_input_images=include_input_images,
            include_output_images=include_output_images,
        )
    elif (
        source_mimetype
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        scraped_chunks = scrape_docx(
            file_path=filepath,
            verbose=verbose,
            include_output_images=include_output_images,
        )
    elif (
        source_mimetype
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        scraped_chunks = scrape_pptx(
            file_path=filepath,
            verbose=verbose,
            include_output_images=include_output_images,
        )
    elif source_mimetype.startswith("image/"):
        scraped_chunks = scrape_image(file_path=filepath)
    elif (
        source_mimetype.startswith("application/vnd.ms-excel")
        or source_mimetype
        == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    ):
        scraped_chunks = scrape_spreadsheet(
            file_path=filepath, source_type=source_mimetype
        )
    elif source_mimetype == "application/x-ipynb+json":
        scraped_chunks = scrape_ipynb(
            file_path=filepath,
            verbose=verbose,
            include_output_images=include_output_images,
        )
    elif (
        source_mimetype == "application/zip"
        or source_mimetype == "application/x-zip-compressed"
    ):
        scraped_chunks = scrape_zip(
            file_path=filepath,
            verbose=verbose,
            openai_client=openai_client,
            include_input_images=include_input_images,
            include_output_images=include_output_images,
        )
    elif source_mimetype.startswith("video/"):
        scraped_chunks = scrape_video(
            file_path=filepath,
            verbose=verbose,
            include_output_images=include_output_images,
        )
    elif source_mimetype.startswith("audio/"):
        scraped_chunks = scrape_audio(file_path=filepath, verbose=verbose)
    elif source_mimetype.startswith("text/html"):
        scraped_chunks = scrape_html(
            file_path=filepath,
            verbose=verbose,
            include_output_images=include_output_images,
        )
    elif source_mimetype.startswith("text/"):
        scraped_chunks = scrape_plaintext(file_path=filepath)
    else:
        try:
            scraped_chunks = scrape_plaintext(file_path=filepath)
        except Exception as e:
            if verbose:
                print(f"[thepipe] Error extracting from {filepath}: {e}")
    if verbose:
        if scraped_chunks:
            print(f"[thepipe] Extracted from {filepath}")
        else:
            print(f"[thepipe] No content extracted from {filepath}")
    if chunking_method:
        scraped_chunks = chunking_method(scraped_chunks)
    return scraped_chunks


def scrape_html(
    file_path: str,
    verbose: bool = False,
    include_output_images: bool = True,
) -> List[Chunk]:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
        html_content = file.read()
    markdown_content = markdownify.markdownify(html_content, heading_style="ATX")
    images = get_images_from_markdown(html_content) if include_output_images else []
    return [Chunk(path=file_path, text=markdown_content, images=images)]


def scrape_plaintext(file_path: str) -> List[Chunk]:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
        text = file.read()
    return [Chunk(path=file_path, text=text)]


def scrape_directory(
    dir_path: str,
    inclusion_pattern: Optional[str] = None,
    verbose: bool = False,
    openai_client: Optional[OpenAI] = None,
    model: str = DEFAULT_AI_MODEL,
    include_input_images: bool = True,
    include_output_images: bool = True,
    _root_dir: Optional[str] = None,
    _visited_dirs: Optional[set[str]] = None,
) -> List[Chunk]:
    """
    inclusion_pattern: Optional regex string; only files whose path matches this pattern will be scraped.
    By default, ignores all files in baked-in constants FOLDERS_TO_IGNORE and FILES_TO_IGNORE.
    """
    # compile the include pattern once
    pattern = re.compile(inclusion_pattern) if inclusion_pattern else None
    extraction: List[Chunk] = []
    canonical_root = os.path.realpath(_root_dir or dir_path)
    current_dir = os.path.realpath(dir_path)
    visited_dirs = _visited_dirs if _visited_dirs is not None else set()

    if not _is_within_directory(canonical_root, current_dir):
        if verbose:
            print(f"[thepipe] Skipping path outside root: {dir_path}")
        return extraction

    if current_dir in visited_dirs:
        if verbose:
            print(f"[thepipe] Skipping already visited directory: {current_dir}")
        return extraction
    visited_dirs.add(current_dir)

    try:
        for entry in os.scandir(dir_path):
            path = entry.path
            resolved_path = os.path.realpath(path)

            if not _is_within_directory(canonical_root, resolved_path):
                if verbose:
                    print(f"[thepipe] Skipping path outside root: {path}")
                continue

            # skip ignored directories
            if entry.is_dir() and any(
                fnmatch.fnmatch(entry.name, pat) for pat in FOLDERS_TO_IGNORE
            ):
                if verbose:
                    print(f"[thepipe] Skipping ignored directory: {path}")
                continue

            # skip ignored files
            if entry.is_file() and any(
                fnmatch.fnmatch(entry.name, pat) for pat in FILES_TO_IGNORE
            ):
                if verbose:
                    print(f"[thepipe] Skipping ignored file: {path}")
                continue

            if entry.is_file():
                # if include_pattern is set, skip files that don't match
                if pattern and not pattern.search(path):
                    if verbose:
                        print(f"[thepipe] Skipping non-matching file: {path}")
                    continue

                if verbose:
                    print(f"[thepipe] Scraping file: {resolved_path}")
                extraction += scrape_file(
                    filepath=resolved_path,
                    verbose=verbose,
                    openai_client=openai_client,
                    model=model,
                    include_input_images=include_input_images,
                    include_output_images=include_output_images,
                )

            elif entry.is_dir():
                # recurse into subdirectory
                if verbose:
                    print(f"[thepipe] Entering directory: {resolved_path}")
                extraction += scrape_directory(
                    dir_path=resolved_path,
                    inclusion_pattern=inclusion_pattern,
                    verbose=verbose,
                    openai_client=openai_client,
                    model=model,
                    include_input_images=include_input_images,
                    include_output_images=include_output_images,
                    _root_dir=canonical_root,
                    _visited_dirs=visited_dirs,
                )
    except PermissionError as e:
        if verbose:
            print(f"[thepipe] Skipping {dir_path} (permission denied): {e}")

    return extraction


def scrape_zip(
    file_path: str,
    inclusion_pattern: Optional[str] = None,
    verbose: bool = False,
    openai_client: Optional[OpenAI] = None,
    include_input_images: bool = True,
    include_output_images: bool = True,
) -> List[Chunk]:
    chunks = []
    with tempfile.TemporaryDirectory() as temp_dir:
        with zipfile.ZipFile(file_path, "r") as zip_ref:
            zip_ref.extractall(temp_dir)
        chunks = scrape_directory(
            dir_path=temp_dir,
            inclusion_pattern=inclusion_pattern,
            verbose=verbose,
            openai_client=openai_client,
            include_input_images=include_input_images,
            include_output_images=include_output_images,
        )
    return chunks


def scrape_pdf(
    file_path: str,
    openai_client: Optional[OpenAI] = None,
    model: str = DEFAULT_AI_MODEL,
    verbose: Optional[bool] = False,
    include_input_images: bool = True,
    include_output_images: bool = True,
    image_scale: float = 1.0,
) -> List[Chunk]:
    chunks: List[Chunk] = []

    # Branch 1 â€“ VLM path (OpenAI client supplied)
    if openai_client is not None:
        with open(file_path, "rb") as fp:
            pdf_bytes = fp.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        num_pages = len(doc)

        if verbose:
            print(
                f"[thepipe] Scraping PDF: {file_path} "
                f"({num_pages} pages) with model {model}"
            )

        # Inner worker â€“ processes one page
        def _process_page(page_num: int) -> Tuple[int, str, Optional[Image.Image]]:
            page = doc[page_num]
            text = page.get_text()  # type: ignore[attr-defined]

            # Build message for the LLM
            msg_content: List[Dict[str, Union[Dict[str, str], str]]] = [
                {
                    "type": "text",
                    "text": f"```\n{text}\n```\n{SCRAPING_PROMPT}",
                }
            ]

            image: Optional[Image.Image] = None
            if include_input_images or include_output_images:
                mat = fitz.Matrix(image_scale, image_scale)
                pix = page.get_pixmap(matrix=mat, alpha=False)  # type: ignore[attr-defined]  # noqa: E501
                image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

                if include_input_images:
                    encoded = make_image_url(image, host_images=HOST_IMAGES)
                    msg_content.append(
                        {
                            "type": "image_url",
                            "image_url": {"url": encoded, "detail": "high"},
                        }
                    )

            messages = cast(
                Iterable[ChatCompletionMessageParam],
                [{"role": "user", "content": msg_content}],
            )

            response = openai_client.chat.completions.create(
                model=model, messages=messages
            )

            llm_response = response.choices[0].message.content
            if not llm_response:
                raise RuntimeError("Empty LLM response.")

            llm_response = llm_response.strip()
            if llm_response.startswith("```markdown"):
                llm_response = llm_response[len("```markdown") :]
            elif llm_response.startswith("```"):
                llm_response = llm_response[len("```") :]
            if llm_response.endswith("```"):
                llm_response = llm_response[: -len("```")]

            return (
                page_num,
                llm_response,
                image if include_output_images else None,
            )

        # Parallel extraction
        max_workers = (os.cpu_count() or 1) * 2
        if verbose:
            print(f"[thepipe] Using {max_workers} threads for PDF extraction")

        page_results: OrderedDict[int, Tuple[str, Optional[Image.Image]]] = (
            OrderedDict()
        )
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = [executor.submit(_process_page, p) for p in range(num_pages)]
            for fut in as_completed(futures):
                pg, txt, img = fut.result()
                page_results[pg] = (txt, img)

        for pg in sorted(page_results):
            txt, img = page_results[pg]
            chunks.append(Chunk(path=file_path, text=txt, images=[img] if img else []))

        return chunks

    # Branch 2 â€“ no OpenAI client â€“ text-only offline mode
    from pymupdf4llm.helpers.pymupdf_rag import to_markdown  # local import

    doc = fitz.open(file_path)
    md_pages = cast(List[Dict[str, Any]], to_markdown(file_path, page_chunks=True))

    for i in range(doc.page_count):
        text = re.sub(r"\n{3,}", "\n\n", md_pages[i]["text"]).strip()

        images: List[Image.Image] = []
        if include_output_images:
            mat = fitz.Matrix(image_scale, image_scale)
            pix = doc[i].get_pixmap(matrix=mat, alpha=False)  # type: ignore[attr-defined]  # noqa: E501
            images.append(Image.frombytes("RGB", [pix.width, pix.height], pix.samples))

        chunks.append(Chunk(path=file_path, text=text, images=images))

    doc.close()
    return chunks


def get_images_from_markdown(text: str) -> List[Image.Image]:
    image_urls = re.findall(r"!\[.*?\]\((.*?)\)", text)
    images = []
    for url in image_urls:
        extension = os.path.splitext(urlparse(url).path)[1]
        if extension not in {".jpg", ".jpeg", ".png"}:
            # ignore incompatible image extractions
            continue

        try:
            response = requests.get(
                url,
                timeout=10,
                headers={"User-Agent": USER_AGENT_STRING},
            )
            response.raise_for_status()
        except Exception:
            continue

        img = Image.open(BytesIO(response.content))
        images.append(img)
    return images


def scrape_image(file_path: str) -> List[Chunk]:
    img = Image.open(file_path)
    img.load()  # needed to close the file
    chunk = Chunk(path=file_path, images=[img])
    return [chunk]


def scrape_spreadsheet(file_path: str, source_type: str) -> List[Chunk]:
    import pandas as pd

    if source_type == "application/vnd.ms-excel":
        df = pd.read_csv(file_path)
    elif (
        source_type
        == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    ):
        df = pd.read_excel(file_path)
    else:
        raise ValueError("Unsupported file format")
    dicts = df.to_dict(orient="records")
    chunks = []
    for i, item in enumerate(dicts):
        # format each row as json along with the row index
        item["row index"] = i
        item_json = json.dumps(item, indent=4)
        chunks.append(Chunk(path=file_path, text=item_json))
    return chunks


def format_timestamp(seconds, chunk_index, chunk_duration):
    # helper function to format the timestamp.
    total_seconds = chunk_index * chunk_duration + seconds
    hours = int(total_seconds // 3600)
    minutes = int((total_seconds % 3600) // 60)
    seconds = total_seconds % 60
    milliseconds = int((seconds - int(seconds)) * 1000)
    return f"{hours:02}:{minutes:02}:{int(seconds):02}.{milliseconds:03}"


def scrape_video(
    file_path: str,
    verbose: bool = False,
    include_output_images: bool = True,
) -> List[Chunk]:
    whisper = _load_whisper()
    from moviepy.editor import VideoFileClip

    # Splits the video into chunks of length MAX_WHISPER_DURATION, extracts
    # one representative frame from the start of each chunk, and then transcribes
    # that chunk.
    model = whisper.load_model("base")
    video = VideoFileClip(file_path)
    num_chunks = math.ceil(video.duration / MAX_WHISPER_DURATION)
    chunks = []

    try:
        for i in range(num_chunks):
            # Calculate the start and end time of the chunk
            start_time = i * MAX_WHISPER_DURATION
            end_time = start_time + MAX_WHISPER_DURATION
            if end_time > video.duration:
                end_time = video.duration

            # Extract a frame from the start of the chunk
            image = None
            if include_output_images:
                frame = video.get_frame(start_time)
                image = Image.fromarray(frame)

            # Save the audio to a temporary .wav file
            with tempfile.NamedTemporaryFile(
                suffix=".wav", delete=False
            ) as temp_audio_file:
                audio_path = temp_audio_file.name

            audio = video.subclip(start_time, end_time).audio  # type: ignore[attr-defined]
            transcription = None

            if audio is not None:
                audio.write_audiofile(audio_path, codec="pcm_s16le")
                result = model.transcribe(audio=audio_path, verbose=verbose)

                # Format transcription with timestamps
                formatted_transcription = []
                for segment in cast(List[Dict[str, Any]], result["segments"]):
                    seg_start = format_timestamp(
                        segment["start"], i, MAX_WHISPER_DURATION
                    )
                    seg_end = format_timestamp(segment["end"], i, MAX_WHISPER_DURATION)
                    formatted_transcription.append(
                        f"[{seg_start} --> {seg_end}]  {segment['text']}"
                    )

                transcription = "\n".join(formatted_transcription)
                os.remove(audio_path)

            # Only add chunks if there is either text or images
            if transcription or image:
                chunks.append(
                    Chunk(
                        path=file_path,
                        text=transcription if transcription else None,
                        images=[image] if image else [],
                    )
                )
    finally:
        video.close()

    return chunks


def scrape_audio(file_path: str, verbose: bool = False) -> List[Chunk]:
    whisper = _load_whisper()

    model = whisper.load_model("base")
    result = model.transcribe(audio=file_path, verbose=verbose)
    segments = cast(List[Dict[str, Any]], result.get("segments", []))

    transcript: List[str] = []
    for segment in segments:
        start = format_timestamp(segment["start"], 0, 0)
        end = format_timestamp(segment["end"], 0, 0)
        if segment["text"].strip():
            transcript.append(f"[{start} --> {end}]  {segment['text']}")
    # join the formatted transcription into a single string
    transcript_text = "\n".join(transcript)
    return [Chunk(path=file_path, text=transcript_text)]


def scrape_docx(
    file_path: str,
    verbose: bool = False,
    include_output_images: bool = True,
) -> List[Chunk]:
    from docx import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table, _Cell
    from docx.text.paragraph import Paragraph
    import csv
    import io

    # helper function to iterate through blocks in the document
    def iter_block_items(parent):
        if parent.__class__.__name__ == "Document":
            parent_elm = parent.element.body
        elif parent.__class__.__name__ == "_Cell":
            parent_elm = parent._tc
        else:
            raise ValueError("Unsupported parent type")
        # iterate through each child element in the parent element
        for child in parent_elm.iterchildren():
            child_elem_class_name = child.__class__.__name__
            if verbose:
                print(f"[thepipe] Found element in docx: {child_elem_class_name}")
            if child_elem_class_name == "CT_P":
                yield Paragraph(child, parent)
            elif child_elem_class_name == "CT_Tbl":
                yield Table(child, parent)

    # helper function to read tables in the document
    def read_docx_tables(tab):
        vf = StringIO()
        writer = csv.writer(vf)
        for row in tab.rows:
            writer.writerow(cell.text for cell in row.cells)
        vf.seek(0)
        return vf.getvalue()

    # read the document
    document = Document(file_path)
    chunks = []
    image_counter = 0

    # Define namespaces
    nsmap = {
        "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
        "pic": "http://schemas.openxmlformats.org/drawingml/2006/picture",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }

    try:
        # scrape each block in the document to create chunks
        # A block can be a paragraph, table, or image
        for block in iter_block_items(document):
            block_texts = []
            block_images = []
            if isinstance(block, Paragraph):
                block_texts.append(block.text)
                # "runs" are the smallest units in a paragraph
                for run in block.runs:
                    if "pic:pic" in run.element.xml and include_output_images:
                        # extract images from the paragraph
                        for pic in run.element.findall(".//pic:pic", nsmap):
                            cNvPr = pic.find(".//pic:cNvPr", nsmap)
                            name_attr = (
                                cNvPr.get("name")
                                if cNvPr is not None
                                else f"image_{image_counter}"
                            )
                            blip = pic.find(".//a:blip", nsmap)
                            if blip is not None:
                                embed_attr = blip.get(
                                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                                )
                                if embed_attr:
                                    image_part = document.part.related_parts[embed_attr]
                                    image_data = BytesIO(image_part._blob)
                                    image = Image.open(image_data)
                                    image.load()
                                    block_images.append(image)
                                    image_counter += 1
            elif isinstance(block, Table):
                table_text = read_docx_tables(block)
                block_texts.append(table_text)
            if block_texts or block_images:
                block_text = "\n".join(block_texts).strip()
                if block_text or block_images:
                    chunks.append(
                        Chunk(path=file_path, text=block_text, images=block_images)
                    )
    except Exception as e:
        raise ValueError(f"Error processing DOCX file {file_path}: {e}")
    return chunks


def scrape_pptx(
    file_path: str,
    verbose: bool = False,
    include_output_images: bool = True,
) -> List[Chunk]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.shapes.picture import Picture
    from pptx.shapes.autoshape import Shape as AutoShape

    prs = Presentation(file_path)
    chunks = []
    # iterate through each slide in the presentation
    for slide in prs.slides:
        slide_texts = []
        slide_images = []
        # iterate through each shape in the slide
        for shape in slide.shapes:
            if shape.has_text_frame:
                auto_shape = cast(AutoShape, shape)
                for paragraph in auto_shape.text_frame.paragraphs:
                    text = paragraph.text
                    if len(slide_texts) == 0:
                        text = "# " + text  # header for first text of a slide
                    slide_texts.append(text)
            # extract images from shapes
            if include_output_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pic = cast(Picture, shape)
                image_data = pic.image.blob
                image = Image.open(BytesIO(image_data))
                slide_images.append(image)
        # add slide to chunks if it has text or images
        if slide_texts or slide_images:
            text = "\n".join(slide_texts).strip()
            if not include_output_images:
                slide_images = []
            chunks.append(Chunk(path=file_path, text=text, images=slide_images))
    # return all chunks
    return chunks


def scrape_ipynb(
    file_path: str,
    verbose: bool = False,
    include_output_images: bool = True,
) -> List[Chunk]:
    with open(file_path, "r", encoding="utf-8") as file:
        notebook = json.load(file)
    chunks = []
    # parse cells in the notebook
    for cell in notebook["cells"]:
        texts = []
        images: List[Image.Image] = []
        cell_type = cell["cell_type"]
        # parse cell content based on type
        if verbose:
            print(f"[thepipe] Scraping cell {cell_type} from {file_path}")
        if cell_type == "markdown":
            text = "".join(cell["source"])
            if include_output_images:
                images = get_images_from_markdown(text)
            texts.append(text)
        elif cell_type == "code":
            source = "".join(cell["source"])
            texts.append(source)
            output_texts = []
            # code cells can have outputs
            if "outputs" in cell:
                for output in cell["outputs"]:
                    if (
                        include_output_images
                        and "data" in output
                        and "image/png" in output["data"]
                    ):
                        image_data = output["data"]["image/png"]
                        image = Image.open(BytesIO(base64.b64decode(image_data)))
                        images.append(image)
                    elif "data" in output and "text/plain" in output["data"]:
                        output_text = "".join(output["data"]["text/plain"])
                        output_texts.append(output_text)
            if output_texts:
                texts.extend(output_texts)
        elif cell_type == "raw":
            text = "".join(cell["source"])
            texts.append(text)
        if texts or images:
            text = "\n".join(texts).strip()
            chunks.append(Chunk(path=file_path, text=text, images=images))
    return chunks
